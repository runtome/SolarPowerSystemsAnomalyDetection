"""Generate presentation: Anomaly Detection in Solar Power Systems."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color palette ---
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_RED = RGBColor(0xE0, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
SUBTLE_BG = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_HEADER_BG = RGBColor(0x00, 0x70, 0xA0)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF4, 0xF8)
SECTION_BG = RGBColor(0x00, 0x50, 0x80)

FIG_DIR = "figures/"


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        if bold_first and ": " in item:
            bold_part, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = bold_part + ": "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return tf


def add_table(slide, left, top, width, height, data, col_widths=None):
    """data: list of lists, first row = header."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
    return table


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)


def title_bar(slide, title, subtitle=None):
    """Dark header bar at top of slide."""
    add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BG)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
                    subtitle, font_size=16, color=LIGHT_GRAY)


def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, SECTION_BG)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                    subtitle, font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4), Inches(3.55), Inches(5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()
    return slide


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "Anomaly Detection in Solar Power Systems",
            font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Using Machine Learning and Deep Learning",
            font_size=28, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(3.5), Inches(4.0), Inches(6), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
            "Site 1 Analysis  |  June - December 2025  |  PyTorch",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Problem Statement")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
            "Goal: Detect anomalies in solar power generation that indicate equipment faults, "
            "panel degradation, inverter failures, or sensor errors.",
            font_size=20, color=DARK_TEXT, bold=True)

add_textbox(slide, Inches(0.6), Inches(2.3), Inches(5.5), Inches(0.4),
            "Why It Matters", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(2.5), [
    "Undetected faults reduce energy output and revenue",
    "Early detection enables preventive maintenance",
    "Reduces system downtime and extends equipment lifespan",
    "Critical for scaling solar infrastructure reliably",
], font_size=17, color=DARK_TEXT)

add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
            "Approach", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.8), Inches(5.8), Inches(3), [
    "Train models to learn the NORMAL relationship between weather conditions and solar generation",
    "When actual output deviates significantly from predicted output, flag as ANOMALY",
    "Use multiple models (ML + DL) and ensemble voting for robust detection",
], font_size=17, color=DARK_TEXT)

# ============================================================
# SLIDE 3: Data Description
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Data Description", "Site 1 Solar Power Plant  |  June - December 2025")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.2), [
    ["Sensor", "Variable", "Unit", "Role", "Resolution", "Raw Rows"],
    ["Generation", "Solar production", "kW", "Target (Y)", "~1 min", "26,226"],
    ["Irradiance", "Solar radiation", "W/m\u00b2", "Feature (X)", "~1 min", "26,177"],
    ["Temperature", "Ambient temp", "\u00b0C", "Feature (X)", "~1 min", "26,195"],
    ["Load", "Power consumption", "kW", "Context", "15 min", "17,367"],
])

add_textbox(slide, Inches(0.6), Inches(4.0), Inches(5), Inches(0.4),
            "Variable Statistics (After Cleaning)", font_size=18, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(0.6), Inches(4.5), Inches(7), Inches(1.8), [
    ["Variable", "Mean", "Std", "Min", "Max"],
    ["Generation (kW)", "96.53", "150.86", "0.00", "638.00"],
    ["Irradiance (W/m\u00b2)", "157.76", "243.96", "0.00", "1,073.00"],
    ["Temperature (\u00b0C)", "28.73", "3.01", "0.00", "37.05"],
])

# ============================================================
# SECTION: EDA
# ============================================================
section_slide("Exploratory Data Analysis", "Understanding the data before modeling")

# ============================================================
# SLIDE 4: Data Preprocessing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Data Preprocessing")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(0.4),
            "Challenge: Different Time Resolutions", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5), [
    "Gen, Irradiance, Temp: ~1-minute intervals",
    "Load: 15-minute intervals",
    "Cannot merge directly - need uniform resolution",
], font_size=16)

add_textbox(slide, Inches(0.6), Inches(3.5), Inches(6), Inches(0.4),
            "Solution Pipeline", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3), [
    "1. Resample all data to uniform 15-min intervals (mean aggregation)",
    "2. Merge 4 variables via outer join on timestamp",
    "3. Forward-fill small gaps (up to 1 hour)",
    "4. Time-based interpolation for remaining short gaps",
    "5. Drop rows with large unrecoverable gaps",
], font_size=16)

add_table(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(1.8), [
    ["Stage", "Rows", "Note"],
    ["Raw merged", "17,664", "Outer join"],
    ["After cleaning", "15,013", "Usable data"],
    ["Dropped", "2,651", "15% loss"],
])

add_image_safe(slide, FIG_DIR + "eda_cell10.png",
               Inches(7), Inches(3.8), width=Inches(5.8))

# ============================================================
# SLIDE 5: EDA - Correlations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: Correlation Analysis")

add_image_safe(slide, FIG_DIR + "eda_cell22.png",
               Inches(0.4), Inches(1.4), width=Inches(5.5))

add_table(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(1.8), [
    ["Feature", "Correlation", "Strength"],
    ["Irradiance (W/m\u00b2)", "0.901", "Very Strong"],
    ["Temperature (\u00b0C)", "0.564", "Moderate"],
    ["Load (kW)", "0.026", "Negligible"],
])

add_textbox(slide, Inches(6.5), Inches(3.5), Inches(6.3), Inches(0.4),
            "Key Insight", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(6.7), Inches(4.0), Inches(6), Inches(3), [
    "Irradiance: r=0.90 - primary driver of generation (physics-based)",
    "Temperature: r=0.56 - moderate, co-varies with sunny hours",
    "Load: r=0.03 - independent consumption pattern",
    "Selected features for modeling: Irradiance + Temperature",
], font_size=16, bold_first=True)

# ============================================================
# SLIDE 6: EDA - Time Series
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: Time Series Overview (6 Months)")

add_image_safe(slide, FIG_DIR + "eda_cell16.png",
               Inches(0.3), Inches(1.3), width=Inches(12.7))

# ============================================================
# SLIDE 7: EDA - Patterns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: Daily Patterns & Scatter Analysis")

add_image_safe(slide, FIG_DIR + "eda_cell25.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "eda_cell23.png",
               Inches(6.6), Inches(1.3), width=Inches(6.5))

# ============================================================
# SLIDE 8: EDA - Key Observations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: Key Observations")

add_image_safe(slide, FIG_DIR + "eda_cell30.png",
               Inches(0.4), Inches(1.4), width=Inches(6))
add_image_safe(slide, FIG_DIR + "eda_cell32.png",
               Inches(0.4), Inches(4.2), width=Inches(6))

add_bullet_list(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5), [
    "1,655 data points: Generation recorded with zero sunlight (potential sensor/logging error)",
    "0 data points: High irradiance with zero generation (no obvious equipment failure detected)",
    "Generation efficiency (daytime): mean = 0.71 kW/(W/m\u00b2)",
    "Solar surplus covers only 10.8% of time",
    "Site relies on grid power 89.2% of the time",
    "Clear daily bell-curve pattern for generation/irradiance",
    "Seasonal variation visible across Jun-Dec period",
], font_size=16, bold_first=False)

# ============================================================
# SECTION: Modeling
# ============================================================
section_slide("Model Training", "Machine Learning & Deep Learning Approaches")

# ============================================================
# SLIDE 9: Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Methodology Overview")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
            "Train/Test Split (Time-Based, No Shuffle)", font_size=20, color=ACCENT_BLUE, bold=True)
add_table(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.4), [
    ["Set", "Rows", "Period", "Ratio"],
    ["Train", "10,509", "Jun 1 - Oct 6, 2025", "70%"],
    ["Test", "4,504", "Oct 6 - Dec 1, 2025", "30%"],
])

add_image_safe(slide, FIG_DIR + "train_cell06.png",
               Inches(6.5), Inches(1.4), width=Inches(6.5))

add_textbox(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.4),
            "6 Models Trained", font_size=20, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.8), [
    ["#", "Model", "Category", "Approach"],
    ["1", "Isolation Forest", "ML, Unsupervised", "Detects statistically isolated points in feature space"],
    ["2", "Random Forest", "ML, Supervised", "Regression prediction; large error = anomaly"],
    ["3", "LSTM", "DL, Sequential", "Learns temporal patterns from 1-hour sliding windows"],
    ["4", "CNN-LSTM", "DL, Hybrid", "CNN extracts local features + LSTM captures time dynamics"],
    ["5", "LSTM Autoencoder", "DL, Reconstruction", "Compresses & reconstructs normal patterns"],
    ["6", "Transformer", "DL, Self-Attention", "Multi-head attention captures all-timestep dependencies"],
])

# ============================================================
# SLIDE 10: DL Model Architectures
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Deep Learning Model Architectures", "PyTorch  |  Input: 4 timesteps x 3 features  |  Output: 3 features")

# LSTM box
add_shape_bg(slide, Inches(0.4), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(0.5), Inches(1.45), Inches(2.8), Inches(0.35),
            "LSTM (53,123 params)", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(1.85), Inches(2.7), Inches(2), [
    "2-layer LSTM (hidden=64)",
    "Dropout 0.2",
    "FC: 64 -> 32 -> 3",
    "Last timestep output",
], font_size=13)

# CNN-LSTM box
add_shape_bg(slide, Inches(3.6), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(3.7), Inches(1.45), Inches(2.8), Inches(0.35),
            "CNN-LSTM (8,819 params)", font_size=16, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(3.8), Inches(1.85), Inches(2.7), Inches(2), [
    "Conv1D: 3->32->16 filters",
    "1-layer LSTM (hidden=32)",
    "FC: 32 -> 32 -> 3",
    "Most parameter-efficient",
], font_size=13)

# Autoencoder box
add_shape_bg(slide, Inches(6.8), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(6.9), Inches(1.45), Inches(2.8), Inches(0.35),
            "LSTM Autoencoder (16,611)", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(1.85), Inches(2.7), Inches(2), [
    "Encoder: LSTM 3->32->16",
    "Bottleneck: 16-dim",
    "Decoder: LSTM 16->32->3",
    "Reconstruct input sequence",
], font_size=13)

# Transformer box
add_shape_bg(slide, Inches(10.0), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(10.1), Inches(1.45), Inches(2.8), Inches(0.35),
            "Transformer (~37K params)", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(10.2), Inches(1.85), Inches(2.7), Inches(2), [
    "Input proj + Pos Encoding",
    "2 Encoder layers, 2 heads",
    "d_model=32, ff_dim=64",
    "Global avg pool + FC",
], font_size=13)

add_textbox(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4),
            "Training Configuration", font_size=18, color=ACCENT_BLUE, bold=True)
add_table(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(2.2), [
    ["Parameter", "Value", "Parameter", "Value"],
    ["Optimizer", "Adam (lr=0.001)", "Batch Size", "32"],
    ["LR Scheduler", "ReduceLROnPlateau (factor=0.5)", "Early Stopping", "Patience=15"],
    ["Loss Function", "MSE", "Gradient Clipping", "Max norm=1.0"],
    ["Scaling", "MinMaxScaler (0-1)", "GPU", "NVIDIA RTX 3060 Ti"],
])

# ============================================================
# SLIDE 11: Anomaly Detection Method
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Anomaly Detection Method", "3-Sigma Rule on Prediction Error")

add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(2.5), WHITE)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
            "Prediction-Based (RF, LSTM, CNN-LSTM, Transformer)", font_size=17, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.1), Inches(5), Inches(1.8), [
    "1. Error = |Actual Generation - Predicted Generation|",
    "2. Threshold = Mean(Error) + 3 x Std(Error)",
    "3. If Error > Threshold  -->  ANOMALY",
], font_size=16, bold_first=True)

add_shape_bg(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.5), WHITE)
add_textbox(slide, Inches(6.7), Inches(1.6), Inches(5.8), Inches(0.4),
            "Reconstruction-Based (LSTM Autoencoder)", font_size=17, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(1.8), [
    "1. Reconstruct input sequence through bottleneck",
    "2. Error = |Original Sequence - Reconstructed Sequence|",
    "3. Threshold from training data (3-sigma)",
    "High error = pattern not learned = anomalous",
], font_size=16, bold_first=True)

add_shape_bg(slide, Inches(0.6), Inches(4.3), Inches(12.2), Inches(2.8), WHITE)
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
            "Why 3-Sigma?", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(4.9), Inches(11), Inches(2), [
    "In a normal distribution, 99.7% of data falls within 3 standard deviations",
    "Points beyond 3-sigma are statistically rare (0.3%) and likely represent genuine anomalies",
    "Tunable: 2-sigma = more sensitive (more detections, more false positives) | 4-sigma = more conservative",
    "Consistent threshold applied across all prediction-based models for fair comparison",
], font_size=16)

# ============================================================
# SECTION: Results
# ============================================================
section_slide("Results", "Prediction Performance & Anomaly Detection")

# ============================================================
# SLIDE 12: Training Curves
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Training Loss Curves")

add_image_safe(slide, FIG_DIR + "train_cell22.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell26.png",
               Inches(6.7), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell30.png",
               Inches(0.2), Inches(4.1), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell34.png",
               Inches(6.7), Inches(4.1), width=Inches(6.4))

# ============================================================
# SLIDE 13: Prediction Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Prediction Performance (Test Set)")

add_table(slide, Inches(0.6), Inches(1.5), Inches(8), Inches(2.2), [
    ["Model", "RMSE", "MAE", "R\u00b2", "Parameters"],
    ["Random Forest", "83.37", "34.80", "0.6346", "-"],
    ["LSTM", "32.46", "15.43", "0.9445", "53,123"],
    ["CNN-LSTM", "32.87", "15.45", "0.9431", "8,819"],
    ["Transformer *", "32.05", "13.76", "0.9459", "~37K"],
])

add_textbox(slide, Inches(0.6), Inches(3.9), Inches(8), Inches(0.3),
            "* Best model", font_size=14, color=ACCENT_ORANGE, bold=True)

add_image_safe(slide, FIG_DIR + "train_cell37.png",
               Inches(0.3), Inches(4.3), width=Inches(12.5))

# ============================================================
# SLIDE 14: Prediction vs Actual plots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Prediction vs Actual: LSTM & CNN-LSTM")

add_image_safe(slide, FIG_DIR + "train_cell23.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell27.png",
               Inches(6.7), Inches(1.3), width=Inches(6.4))

# ============================================================
# SLIDE 15: Prediction vs Actual - Transformer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Prediction vs Actual: Transformer (Best Model)")

add_image_safe(slide, FIG_DIR + "train_cell35.png",
               Inches(0.3), Inches(1.3), width=Inches(12.5))

# ============================================================
# SLIDE 16: Anomaly Detection Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Anomaly Detection Results")

add_table(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(2.8), [
    ["Model", "Method", "Anomalies", "% of Test"],
    ["Isolation Forest", "Isolation score", "244", "5.4%"],
    ["Random Forest", "3-sigma MAE", "159", "3.5%"],
    ["LSTM", "3-sigma MAE", "102", "2.3%"],
    ["CNN-LSTM", "3-sigma MAE", "94", "2.1%"],
    ["LSTM Autoencoder", "3-sigma recon.", "265", "5.9%"],
    ["Transformer", "3-sigma MAE", "100", "2.2%"],
])

add_image_safe(slide, FIG_DIR + "train_cell38.png",
               Inches(8.3), Inches(1.5), width=Inches(4.8))

add_bullet_list(slide, Inches(0.8), Inches(4.6), Inches(12), Inches(2.5), [
    "Most conservative: CNN-LSTM (94 anomalies, 2.1%)",
    "Most sensitive: LSTM Autoencoder (265, 5.9%) - reconstruction-based catches more subtle deviations",
    "DL prediction models converge at ~2.1-2.3% - consistent and reliable",
    "Unsupervised methods (IF, AE) detect more anomalies but may include more false positives",
], font_size=16, bold_first=True)

# ============================================================
# SLIDE 17: Isolation Forest & RF anomaly plots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Anomaly Visualization: ML Models")

add_image_safe(slide, FIG_DIR + "train_cell16.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell19.png",
               Inches(6.7), Inches(1.3), width=Inches(6.4))

# ============================================================
# SLIDE 18: Autoencoder anomaly plot
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Anomaly Visualization: LSTM Autoencoder")

add_image_safe(slide, FIG_DIR + "train_cell31.png",
               Inches(0.5), Inches(1.3), width=Inches(12))

# ============================================================
# SLIDE 19: Ensemble
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Ensemble Anomaly Detection (Majority Vote)")

add_table(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.0), [
    ["Agreement Level", "Data Points", "%"],
    ["0 models (normal)", "3,939", "87.5%"],
    ["1 model only", "311", "6.9%"],
    ["2 models", "119", "2.6%"],
    ["3 models (anomaly)", "114", "2.5%"],
    ["4 models (anomaly)", "13", "0.3%"],
    ["5 models (anomaly)", "3", "0.07%"],
    ["6 models (anomaly)", "1", "0.02%"],
])

add_shape_bg(slide, Inches(0.6), Inches(4.8), Inches(5.5), Inches(1.2), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(4.9), Inches(5), Inches(1),
            "Total Ensemble Anomalies: 131 (2.9%)\nRule: Anomaly if >= 3 of 6 models agree",
            font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_image_safe(slide, FIG_DIR + "train_cell41.png",
               Inches(6.5), Inches(1.4), width=Inches(6.5))

# ============================================================
# SLIDE 20: Anomaly Interpretation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "What Do These Anomalies Mean?")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.5), [
    ["Anomaly Pattern", "Possible Cause", "Action Required"],
    ["High irradiance + Low generation", "Panel fault, inverter failure, shading", "Inspect panels & inverters"],
    ["Sudden generation drop", "Equipment trip, grid disconnection", "Check breakers & grid connection"],
    ["Generation with zero irradiance", "Sensor malfunction, data logging error", "Calibrate sensors"],
    ["Gradual efficiency decline", "Panel degradation, dust accumulation", "Schedule cleaning/maintenance"],
    ["Unusual generation spikes", "Sensor calibration drift", "Verify sensor readings"],
])

add_textbox(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
            "Practical Impact", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(2.5), [
    "131 flagged time points in ~2 months of test data (Oct-Dec 2025)",
    "Each represents a 15-minute window where the system behaved abnormally",
    "Operations team can investigate clustered anomalies for root cause analysis",
    "Reduces manual inspection effort by focusing attention on high-confidence alerts",
], font_size=17)

# ============================================================
# SLIDE 21: Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Conclusion")

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(5.5), [
    "1. Deep learning dramatically outperforms traditional ML: R\u00b2 improved from 0.63 (Random Forest) to 0.95 (Transformer) - a 49% reduction in unexplained variance.",
    "2. Transformer achieved the best prediction accuracy (RMSE=32.05, MAE=13.76, R\u00b2=0.9459), while CNN-LSTM is the most efficient (comparable accuracy with 6x fewer parameters).",
    "3. Irradiance is the dominant predictor (94.6% feature importance), confirming the physics-based relationship between solar radiation and power generation.",
    "4. Ensemble voting (majority of 6 models) identified 131 high-confidence anomalies (2.9%) in the test period, providing robust detection with reduced false positives.",
    "5. The approach is interpretable: an anomaly means the system produced significantly different power than expected given current weather conditions.",
    "6. 1,655 data points with generation at zero sunlight warrant investigation as potential sensor/data issues.",
], font_size=17, bold_first=True)

# ============================================================
# SLIDE 22: Future Work
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "Future Work")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(3.0), [
    ["Priority", "Task", "Description"],
    ["High", "Multi-site deployment", "Extend to all 4 sites (currently only Site 1)"],
    ["High", "Threshold tuning", "Calibrate sigma per site using operator feedback"],
    ["Medium", "Anomaly classification", "Categorize type: inverter, panel, sensor, etc."],
    ["Medium", "Real-time pipeline", "Build streaming inference for continuous monitoring"],
    ["Low", "Transfer learning", "Pre-train on Site 1, fine-tune on Sites 2-4"],
    ["Low", "Explainability", "Add SHAP / attention visualization for root cause"],
])

# ============================================================
# SLIDE 23: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
            "Thank You",
            font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4), Inches(3.8), Inches(5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_ORANGE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
            "Anomaly Detection in Solar Power Systems\nUsing Machine Learning and Deep Learning",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = "Solar_Anomaly_Detection_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
