"""Generate Thai presentation: Anomaly Detection in Solar Power Systems."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color palette ---
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_RED = RGBColor(0xE0, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
SUBTLE_BG = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_HEADER_BG = RGBColor(0x00, 0x70, 0xA0)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF4, 0xF8)
SECTION_BG = RGBColor(0x00, 0x50, 0x80)

FIG_DIR = "figures/"
THAI_FONT = "Tahoma"


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name=THAI_FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        if bold_first and ": " in item:
            bold_part, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = bold_part + ": "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = THAI_FONT
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = THAI_FONT
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = THAI_FONT
    return tf


def add_table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = THAI_FONT
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
    return table


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)


def title_bar(slide, title, subtitle=None):
    add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BG)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
                    subtitle, font_size=16, color=LIGHT_GRAY)


def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, SECTION_BG)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                    subtitle, font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4), Inches(3.55), Inches(5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()
    return slide


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "การตรวจจับความผิดปกติในระบบผลิตไฟฟ้าพลังงานแสงอาทิตย์",
            font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "ด้วย Machine Learning และ Deep Learning",
            font_size=28, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(3.5), Inches(4.0), Inches(6), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
            "Site 1  |  มิถุนายน - ธันวาคม 2568  |  PyTorch",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "คำชี้แจงปัญหา (Problem Statement)")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
            "เป้าหมาย: ตรวจจับความผิดปกติในการผลิตไฟฟ้าพลังงานแสงอาทิตย์ "
            "ที่บ่งชี้ถึงความเสียหายของอุปกรณ์ แผงโซลาร์เสื่อมสภาพ อินเวอร์เตอร์ขัดข้อง หรือเซ็นเซอร์ผิดพลาด",
            font_size=20, color=DARK_TEXT, bold=True)

add_textbox(slide, Inches(0.6), Inches(2.3), Inches(5.5), Inches(0.4),
            "ทำไมจึงสำคัญ", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(2.5), [
    "ความผิดปกติที่ไม่ถูกตรวจพบ ทำให้สูญเสียพลังงานและรายได้",
    "การตรวจจับล่วงหน้าช่วยให้ซ่อมบำรุงเชิงป้องกันได้",
    "ลดเวลาหยุดทำงานของระบบ ยืดอายุการใช้งานอุปกรณ์",
    "สำคัญต่อการขยายระบบโซลาร์อย่างมีเสถียรภาพ",
], font_size=17)

add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
            "แนวทาง", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.8), Inches(5.8), Inches(3), [
    "สอนโมเดลให้เรียนรู้ความสัมพันธ์ปกติระหว่างสภาพอากาศกับการผลิตไฟฟ้า",
    "เมื่อค่าจริงเบี่ยงเบนจากค่าทำนายอย่างมาก ให้ตั้งสถานะเป็นความผิดปกติ",
    "ใช้หลายโมเดล (ML + DL) และ Ensemble Voting เพื่อเพิ่มความน่าเชื่อถือ",
], font_size=17)

# ============================================================
# SLIDE 3: Data Description
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "รายละเอียดข้อมูล", "โรงไฟฟ้าพลังงานแสงอาทิตย์ Site 1  |  มิ.ย. - ธ.ค. 2568")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.2), [
    ["เซ็นเซอร์", "ตัวแปร", "หน่วย", "บทบาท", "ความละเอียด", "จำนวนแถว"],
    ["การผลิตไฟฟ้า", "กำลังผลิตไฟฟ้า", "kW", "เป้าหมาย (Y)", "~1 นาที", "26,226"],
    ["ความเข้มแสง", "ความเข้มรังสีแสงอาทิตย์", "W/m\u00b2", "ฟีเจอร์ (X)", "~1 นาที", "26,177"],
    ["อุณหภูมิ", "อุณหภูมิแวดล้อม", "\u00b0C", "ฟีเจอร์ (X)", "~1 นาที", "26,195"],
    ["โหลด", "การใช้ไฟฟ้า", "kW", "บริบท", "15 นาที", "17,367"],
])

add_textbox(slide, Inches(0.6), Inches(4.0), Inches(5), Inches(0.4),
            "สถิติตัวแปร (หลังทำความสะอาดข้อมูล)", font_size=18, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(0.6), Inches(4.5), Inches(7), Inches(1.8), [
    ["ตัวแปร", "ค่าเฉลี่ย", "ส่วนเบี่ยงเบน", "ต่ำสุด", "สูงสุด"],
    ["การผลิต (kW)", "96.53", "150.86", "0.00", "638.00"],
    ["ความเข้มแสง (W/m\u00b2)", "157.76", "243.96", "0.00", "1,073.00"],
    ["อุณหภูมิ (\u00b0C)", "28.73", "3.01", "0.00", "37.05"],
])

# ============================================================
# SECTION: EDA
# ============================================================
section_slide("การวิเคราะห์ข้อมูลเชิงสำรวจ (EDA)", "ทำความเข้าใจข้อมูลก่อนการสร้างโมเดล")

# ============================================================
# SLIDE 4: Data Preprocessing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "การเตรียมข้อมูล (Data Preprocessing)")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(0.4),
            "ปัญหา: ความละเอียดเวลาต่างกัน", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5), [
    "การผลิต, ความเข้มแสง, อุณหภูมิ: ทุก ~1 นาที",
    "โหลด: ทุก 15 นาที",
    "ไม่สามารถรวมข้อมูลโดยตรงได้ ต้องปรับความละเอียดให้เท่ากัน",
], font_size=16)

add_textbox(slide, Inches(0.6), Inches(3.5), Inches(6), Inches(0.4),
            "ขั้นตอนการแก้ปัญหา", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3), [
    "1. Resample ข้อมูลทั้งหมดเป็น 15 นาที (ใช้ค่าเฉลี่ย)",
    "2. รวม 4 ตัวแปร ด้วย Outer Join ตาม timestamp",
    "3. Forward-fill ช่องว่างเล็กๆ (ไม่เกิน 1 ชั่วโมง)",
    "4. Interpolation ตามเวลาสำหรับช่องว่างที่เหลือ",
    "5. ลบแถวที่มีช่องว่างขนาดใหญ่ไม่สามารถกู้คืนได้",
], font_size=16)

add_table(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(1.8), [
    ["ขั้นตอน", "จำนวนแถว", "หมายเหตุ"],
    ["หลังรวมข้อมูล", "17,664", "Outer join"],
    ["หลังทำความสะอาด", "15,013", "ข้อมูลพร้อมใช้"],
    ["ถูกลบ", "2,651", "สูญเสีย 15%"],
])

add_image_safe(slide, FIG_DIR + "eda_cell10.png",
               Inches(7), Inches(3.8), width=Inches(5.8))

# ============================================================
# SLIDE 5: EDA - Correlations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: การวิเคราะห์สหสัมพันธ์ (Correlation)")

add_image_safe(slide, FIG_DIR + "eda_cell22.png",
               Inches(0.4), Inches(1.4), width=Inches(5.5))

add_table(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(1.8), [
    ["ฟีเจอร์", "สหสัมพันธ์", "ระดับ"],
    ["ความเข้มแสง (W/m\u00b2)", "0.901", "สูงมาก"],
    ["อุณหภูมิ (\u00b0C)", "0.564", "ปานกลาง"],
    ["โหลด (kW)", "0.026", "ไม่มีนัยสำคัญ"],
])

add_textbox(slide, Inches(6.5), Inches(3.5), Inches(6.3), Inches(0.4),
            "ข้อค้นพบสำคัญ", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(6.7), Inches(4.0), Inches(6), Inches(3), [
    "ความเข้มแสง: r=0.90 - ตัวขับเคลื่อนหลักของการผลิตไฟฟ้า (ตามหลักฟิสิกส์)",
    "อุณหภูมิ: r=0.56 - สัมพันธ์ปานกลาง เพราะแดดจัดมักร้อนด้วย",
    "โหลด: r=0.03 - รูปแบบการใช้ไฟแยกอิสระจากการผลิต",
    "เลือกฟีเจอร์สำหรับโมเดล: ความเข้มแสง + อุณหภูมิ",
], font_size=16, bold_first=True)

# ============================================================
# SLIDE 6: EDA - Time Series
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: ภาพรวมอนุกรมเวลา (6 เดือน)")

add_image_safe(slide, FIG_DIR + "eda_cell16.png",
               Inches(0.3), Inches(1.3), width=Inches(12.7))

# ============================================================
# SLIDE 7: EDA - Patterns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: รูปแบบรายวัน และ การกระจายข้อมูล")

add_image_safe(slide, FIG_DIR + "eda_cell25.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "eda_cell23.png",
               Inches(6.6), Inches(1.3), width=Inches(6.5))

# ============================================================
# SLIDE 8: EDA - Key Observations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "EDA: ข้อสังเกตสำคัญ")

add_image_safe(slide, FIG_DIR + "eda_cell30.png",
               Inches(0.4), Inches(1.4), width=Inches(6))
add_image_safe(slide, FIG_DIR + "eda_cell32.png",
               Inches(0.4), Inches(4.2), width=Inches(6))

add_bullet_list(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5), [
    "พบ 1,655 จุดข้อมูล: มีการผลิตไฟฟ้าทั้งที่ไม่มีแสงแดด (อาจเป็นปัญหาเซ็นเซอร์)",
    "พบ 0 จุดข้อมูล: แสงแดดจัดแต่ไม่มีการผลิต (ไม่พบความเสียหายอุปกรณ์ชัดเจน)",
    "ประสิทธิภาพการผลิต (ช่วงกลางวัน): ค่าเฉลี่ย = 0.71 kW/(W/m\u00b2)",
    "พลังงานแสงอาทิตย์เกินโหลดเพียง 10.8% ของเวลา",
    "ไซต์พึ่งพาไฟฟ้าจากกริด 89.2% ของเวลา",
    "รูปแบบระฆังคว่ำชัดเจนสำหรับการผลิต/ความเข้มแสงรายวัน",
    "เห็นการเปลี่ยนแปลงตามฤดูกาลตลอด มิ.ย. - ธ.ค.",
], font_size=16)

# ============================================================
# SECTION: Modeling
# ============================================================
section_slide("การฝึกสอนโมเดล", "แนวทาง Machine Learning และ Deep Learning")

# ============================================================
# SLIDE 9: Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "ภาพรวมวิธีการ (Methodology)")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
            "การแบ่งข้อมูล Train/Test (ตามลำดับเวลา ไม่สุ่ม)", font_size=20, color=ACCENT_BLUE, bold=True)
add_table(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.4), [
    ["ชุดข้อมูล", "จำนวนแถว", "ช่วงเวลา", "สัดส่วน"],
    ["Train", "10,509", "1 มิ.ย. - 6 ต.ค. 2568", "70%"],
    ["Test", "4,504", "6 ต.ค. - 1 ธ.ค. 2568", "30%"],
])

add_image_safe(slide, FIG_DIR + "train_cell06.png",
               Inches(6.5), Inches(1.4), width=Inches(6.5))

add_textbox(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.4),
            "โมเดลที่ฝึกทั้งหมด 6 โมเดล", font_size=20, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.8), [
    ["#", "โมเดล", "ประเภท", "แนวทาง"],
    ["1", "Isolation Forest", "ML, ไม่มีผู้สอน", "ค้นหาจุดข้อมูลที่แยกตัวออกจากกลุ่มได้ง่าย"],
    ["2", "Random Forest", "ML, มีผู้สอน", "ทำนายการผลิต; ค่าผิดพลาดมาก = ความผิดปกติ"],
    ["3", "LSTM", "DL, อนุกรมเวลา", "เรียนรู้รูปแบบเวลาจากหน้าต่าง 1 ชั่วโมง"],
    ["4", "CNN-LSTM", "DL, ลูกผสม", "CNN ดึงฟีเจอร์เฉพาะที่ + LSTM จับพลวัตเวลา"],
    ["5", "LSTM Autoencoder", "DL, การสร้างใหม่", "บีบอัดและสร้างรูปแบบปกติใหม่"],
    ["6", "Transformer", "DL, Self-Attention", "กลไก Attention จับความสัมพันธ์ข้ามทุก timestep"],
])

# ============================================================
# SLIDE 10: DL Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "สถาปัตยกรรมโมเดล Deep Learning",
          "PyTorch  |  Input: 4 timesteps x 3 features  |  Output: 3 features")

add_shape_bg(slide, Inches(0.4), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(0.5), Inches(1.45), Inches(2.8), Inches(0.35),
            "LSTM (53,123 params)", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(1.85), Inches(2.7), Inches(2), [
    "LSTM 2 ชั้น (hidden=64)",
    "Dropout 0.2",
    "FC: 64 -> 32 -> 3",
    "ใช้ output timestep สุดท้าย",
], font_size=13)

add_shape_bg(slide, Inches(3.6), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(3.7), Inches(1.45), Inches(2.8), Inches(0.35),
            "CNN-LSTM (8,819 params)", font_size=16, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(3.8), Inches(1.85), Inches(2.7), Inches(2), [
    "Conv1D: 3->32->16 filters",
    "LSTM 1 ชั้น (hidden=32)",
    "FC: 32 -> 32 -> 3",
    "ประหยัดพารามิเตอร์ที่สุด",
], font_size=13)

add_shape_bg(slide, Inches(6.8), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(6.9), Inches(1.45), Inches(2.8), Inches(0.35),
            "LSTM Autoencoder (16,611)", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(1.85), Inches(2.7), Inches(2), [
    "Encoder: LSTM 3->32->16",
    "Bottleneck: 16 มิติ",
    "Decoder: LSTM 16->32->3",
    "สร้างลำดับ input ใหม่",
], font_size=13)

add_shape_bg(slide, Inches(10.0), Inches(1.4), Inches(3), Inches(2.7), WHITE)
add_textbox(slide, Inches(10.1), Inches(1.45), Inches(2.8), Inches(0.35),
            "Transformer (~37K params)", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(10.2), Inches(1.85), Inches(2.7), Inches(2), [
    "Input proj + Pos Encoding",
    "Encoder 2 ชั้น, 2 heads",
    "d_model=32, ff_dim=64",
    "Global avg pool + FC",
], font_size=13)

add_textbox(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4),
            "การตั้งค่าการฝึกสอน", font_size=18, color=ACCENT_BLUE, bold=True)
add_table(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(2.2), [
    ["พารามิเตอร์", "ค่า", "พารามิเตอร์", "ค่า"],
    ["Optimizer", "Adam (lr=0.001)", "Batch Size", "32"],
    ["LR Scheduler", "ReduceLROnPlateau (factor=0.5)", "Early Stopping", "Patience=15"],
    ["Loss Function", "MSE", "Gradient Clipping", "Max norm=1.0"],
    ["Scaling", "MinMaxScaler (0-1)", "GPU", "NVIDIA RTX 3060 Ti"],
])

# ============================================================
# SLIDE 11: Anomaly Detection Method
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "วิธีการตรวจจับความผิดปกติ", "กฎ 3-Sigma บนค่าผิดพลาดจากการทำนาย")

add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(2.5), WHITE)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
            "แบบทำนาย (RF, LSTM, CNN-LSTM, Transformer)", font_size=17, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.1), Inches(5), Inches(1.8), [
    "1. ค่าผิดพลาด = |ค่าจริง - ค่าทำนาย|",
    "2. เกณฑ์ = ค่าเฉลี่ย(ผิดพลาด) + 3 x ส่วนเบี่ยงเบน(ผิดพลาด)",
    "3. ถ้าค่าผิดพลาด > เกณฑ์  -->  ผิดปกติ",
], font_size=16, bold_first=True)

add_shape_bg(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.5), WHITE)
add_textbox(slide, Inches(6.7), Inches(1.6), Inches(5.8), Inches(0.4),
            "แบบสร้างใหม่ (LSTM Autoencoder)", font_size=17, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(6.9), Inches(2.1), Inches(5.8), Inches(1.8), [
    "1. สร้างลำดับ input ใหม่ผ่าน bottleneck",
    "2. ค่าผิดพลาด = |ลำดับจริง - ลำดับที่สร้างใหม่|",
    "3. เกณฑ์จากข้อมูล training (3-sigma)",
    "ค่าผิดพลาดสูง = รูปแบบที่โมเดลไม่เคยเรียนรู้ = ผิดปกติ",
], font_size=16, bold_first=True)

add_shape_bg(slide, Inches(0.6), Inches(4.3), Inches(12.2), Inches(2.8), WHITE)
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
            "ทำไมต้อง 3-Sigma?", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(4.9), Inches(11), Inches(2), [
    "ในการแจกแจงปกติ 99.7% ของข้อมูลอยู่ภายใน 3 ส่วนเบี่ยงเบนมาตรฐาน",
    "จุดที่เกิน 3-sigma มีความหายาก (0.3%) และน่าจะเป็นความผิดปกติจริง",
    "ปรับได้: 2-sigma = ไวกว่า (ตรวจพบมากขึ้น อาจมี false positive มากขึ้น) | 4-sigma = อนุรักษ์นิยมกว่า",
    "ใช้เกณฑ์เดียวกันทุกโมเดลเพื่อการเปรียบเทียบที่เป็นธรรม",
], font_size=16)

# ============================================================
# SECTION: Results
# ============================================================
section_slide("ผลลัพธ์", "ประสิทธิภาพการทำนาย และ การตรวจจับความผิดปกติ")

# ============================================================
# SLIDE 12: Training Curves
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "กราฟ Training Loss")

add_image_safe(slide, FIG_DIR + "train_cell22.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell26.png",
               Inches(6.7), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell30.png",
               Inches(0.2), Inches(4.1), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell34.png",
               Inches(6.7), Inches(4.1), width=Inches(6.4))

# ============================================================
# SLIDE 13: Prediction Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "ผลการทำนาย (Test Set)")

add_table(slide, Inches(0.6), Inches(1.5), Inches(8), Inches(2.2), [
    ["โมเดล", "RMSE", "MAE", "R\u00b2", "พารามิเตอร์"],
    ["Random Forest", "83.37", "34.80", "0.6346", "-"],
    ["LSTM", "32.46", "15.43", "0.9445", "53,123"],
    ["CNN-LSTM", "32.87", "15.45", "0.9431", "8,819"],
    ["Transformer *", "32.05", "13.76", "0.9459", "~37K"],
])

add_textbox(slide, Inches(0.6), Inches(3.9), Inches(8), Inches(0.3),
            "* โมเดลที่ดีที่สุด", font_size=14, color=ACCENT_ORANGE, bold=True)

add_image_safe(slide, FIG_DIR + "train_cell37.png",
               Inches(0.3), Inches(4.3), width=Inches(12.5))

# ============================================================
# SLIDE 14: Prediction vs Actual plots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "ค่าทำนาย vs ค่าจริง: LSTM และ CNN-LSTM")

add_image_safe(slide, FIG_DIR + "train_cell23.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell27.png",
               Inches(6.7), Inches(1.3), width=Inches(6.4))

# ============================================================
# SLIDE 15: Transformer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "ค่าทำนาย vs ค่าจริง: Transformer (โมเดลดีที่สุด)")

add_image_safe(slide, FIG_DIR + "train_cell35.png",
               Inches(0.3), Inches(1.3), width=Inches(12.5))

# ============================================================
# SLIDE 16: Anomaly Detection Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "ผลการตรวจจับความผิดปกติ")

add_table(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(2.8), [
    ["โมเดล", "วิธีการ", "ผิดปกติ", "% ของ Test"],
    ["Isolation Forest", "Isolation score", "244", "5.4%"],
    ["Random Forest", "3-sigma MAE", "159", "3.5%"],
    ["LSTM", "3-sigma MAE", "102", "2.3%"],
    ["CNN-LSTM", "3-sigma MAE", "94", "2.1%"],
    ["LSTM Autoencoder", "3-sigma recon.", "265", "5.9%"],
    ["Transformer", "3-sigma MAE", "100", "2.2%"],
])

add_image_safe(slide, FIG_DIR + "train_cell38.png",
               Inches(8.3), Inches(1.5), width=Inches(4.8))

add_bullet_list(slide, Inches(0.8), Inches(4.6), Inches(12), Inches(2.5), [
    "อนุรักษ์นิยมที่สุด: CNN-LSTM (94 จุด, 2.1%)",
    "ไวที่สุด: LSTM Autoencoder (265 จุด, 5.9%) - แบบสร้างใหม่จับความเบี่ยงเบนละเอียดกว่า",
    "โมเดล DL แบบทำนาย ลู่เข้าที่ ~2.1-2.3% - สอดคล้องและเชื่อถือได้",
    "วิธี Unsupervised (IF, AE) ตรวจพบมากกว่า แต่อาจมี false positive มากขึ้น",
], font_size=16, bold_first=True)

# ============================================================
# SLIDE 17: ML Anomaly Plots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "การแสดงผลความผิดปกติ: โมเดล ML")

add_image_safe(slide, FIG_DIR + "train_cell16.png",
               Inches(0.2), Inches(1.3), width=Inches(6.4))
add_image_safe(slide, FIG_DIR + "train_cell19.png",
               Inches(6.7), Inches(1.3), width=Inches(6.4))

# ============================================================
# SLIDE 18: Autoencoder
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "การแสดงผลความผิดปกติ: LSTM Autoencoder")

add_image_safe(slide, FIG_DIR + "train_cell31.png",
               Inches(0.5), Inches(1.3), width=Inches(12))

# ============================================================
# SLIDE 19: Ensemble
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "การตรวจจับแบบ Ensemble (Majority Vote)")

add_table(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.0), [
    ["ระดับความเห็นพ้อง", "จำนวนจุด", "%"],
    ["0 โมเดล (ปกติ)", "3,939", "87.5%"],
    ["1 โมเดล", "311", "6.9%"],
    ["2 โมเดล", "119", "2.6%"],
    ["3 โมเดล (ผิดปกติ)", "114", "2.5%"],
    ["4 โมเดล (ผิดปกติ)", "13", "0.3%"],
    ["5 โมเดล (ผิดปกติ)", "3", "0.07%"],
    ["6 โมเดล (ผิดปกติ)", "1", "0.02%"],
])

add_shape_bg(slide, Inches(0.6), Inches(4.8), Inches(5.5), Inches(1.2), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(4.9), Inches(5), Inches(1),
            "ผิดปกติรวม Ensemble: 131 จุด (2.9%)\nกฎ: ผิดปกติเมื่อ >= 3 ใน 6 โมเดลเห็นพ้อง",
            font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_image_safe(slide, FIG_DIR + "train_cell41.png",
               Inches(6.5), Inches(1.4), width=Inches(6.5))

# ============================================================
# SLIDE 20: Anomaly Interpretation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "ความผิดปกติเหล่านี้หมายถึงอะไร?")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.5), [
    ["รูปแบบความผิดปกติ", "สาเหตุที่เป็นไปได้", "สิ่งที่ต้องดำเนินการ"],
    ["แสงแดดจัด + ผลิตไฟน้อย", "แผงโซลาร์เสีย, อินเวอร์เตอร์ขัดข้อง, เงาบัง", "ตรวจสอบแผงและอินเวอร์เตอร์"],
    ["การผลิตตกฮวบกะทันหัน", "อุปกรณ์ตัด, การเชื่อมต่อกริดขาด", "ตรวจเบรกเกอร์และการเชื่อมต่อกริด"],
    ["ผลิตไฟทั้งที่ไม่มีแสงแดด", "เซ็นเซอร์ทำงานผิดปกติ, ข้อผิดพลาดในการบันทึกข้อมูล", "สอบเทียบเซ็นเซอร์"],
    ["ประสิทธิภาพลดลงอย่างค่อยเป็นค่อยไป", "แผงโซลาร์เสื่อมสภาพ, ฝุ่นสะสม", "กำหนดการทำความสะอาด/บำรุงรักษา"],
    ["กำลังผลิตพุ่งผิดปกติ", "เซ็นเซอร์เลื่อนค่า", "ตรวจสอบค่าเซ็นเซอร์"],
])

add_textbox(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
            "ผลกระทบเชิงปฏิบัติ", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(2.5), [
    "ตรวจพบ 131 จุดเวลาที่ถูกตั้งค่าสถานะ ใน ~2 เดือนของข้อมูลทดสอบ (ต.ค.-ธ.ค. 2568)",
    "แต่ละจุดแทนหน้าต่าง 15 นาที ที่ระบบทำงานผิดปกติ",
    "ทีมปฏิบัติการสามารถตรวจสอบกลุ่มความผิดปกติเพื่อวิเคราะห์สาเหตุราก",
    "ลดภาระการตรวจสอบด้วยตา โดยเน้นเฉพาะการแจ้งเตือนที่มีความน่าเชื่อถือสูง",
], font_size=17)

# ============================================================
# SLIDE 21: Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "สรุปผล")

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(5.5), [
    "1. Deep Learning เหนือกว่า ML แบบดั้งเดิมอย่างชัดเจน: R\u00b2 เพิ่มจาก 0.63 (Random Forest) เป็น 0.95 (Transformer) - ลดความแปรปรวนที่อธิบายไม่ได้ 49%",
    "2. Transformer ทำนายได้แม่นยำที่สุด (RMSE=32.05, MAE=13.76, R\u00b2=0.9459) ขณะที่ CNN-LSTM ประหยัดที่สุด (ความแม่นยำใกล้เคียงด้วยพารามิเตอร์น้อยกว่า 6 เท่า)",
    "3. ความเข้มแสงเป็นตัวทำนายหลัก (ความสำคัญ 94.6%) ยืนยันความสัมพันธ์ทางฟิสิกส์ระหว่างรังสีแสงอาทิตย์กับการผลิตไฟฟ้า",
    "4. Ensemble Voting (เสียงข้างมากจาก 6 โมเดล) ระบุความผิดปกติที่มีความน่าเชื่อถือสูง 131 จุด (2.9%) ในช่วงทดสอบ ลด false positive",
    "5. แนวทางนี้ตีความได้ง่าย: ความผิดปกติหมายถึง ระบบผลิตไฟฟ้าแตกต่างจากที่คาดอย่างมาก เมื่อพิจารณาสภาพอากาศปัจจุบัน",
    "6. พบ 1,655 จุดข้อมูลที่ผลิตไฟฟ้าตอนไม่มีแสงแดด ควรตรวจสอบปัญหาเซ็นเซอร์/ข้อมูล",
], font_size=17, bold_first=True)

# ============================================================
# SLIDE 22: Future Work
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "แผนงานในอนาคต")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(3.0), [
    ["ลำดับ", "งาน", "รายละเอียด"],
    ["สูง", "ขยายไปทุกไซต์", "ขยายไปยังทั้ง 4 ไซต์ (ปัจจุบันมีเฉพาะ Site 1)"],
    ["สูง", "ปรับเกณฑ์", "ปรับค่า sigma ต่อไซต์ตาม feedback จากผู้ปฏิบัติงาน"],
    ["กลาง", "จำแนกประเภทความผิดปกติ", "แยกประเภท: อินเวอร์เตอร์, แผง, เซ็นเซอร์ ฯลฯ"],
    ["กลาง", "ระบบ Real-time", "สร้าง streaming inference สำหรับตรวจสอบต่อเนื่อง"],
    ["ต่ำ", "Transfer Learning", "ฝึกล่วงหน้าบน Site 1, fine-tune บน Sites 2-4"],
    ["ต่ำ", "Explainability", "เพิ่ม SHAP / Attention visualization เพื่อวิเคราะห์สาเหตุ"],
])

# ============================================================
# SLIDE 23: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
            "ขอบคุณครับ/ค่ะ",
            font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4), Inches(3.5), Inches(5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_ORANGE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(1),
            "การตรวจจับความผิดปกติในระบบผลิตไฟฟ้าพลังงานแสงอาทิตย์\nด้วย Machine Learning และ Deep Learning",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = "Solar_Anomaly_Detection_Presentation_TH.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
