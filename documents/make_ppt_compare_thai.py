"""Generate Thai comparison presentation: Variable Timestep Experiments."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color palette (same as original) ---
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_RED = RGBColor(0xE0, 0x40, 0x40)
ACCENT_PURPLE = RGBColor(0x9C, 0x27, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
SUBTLE_BG = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_HEADER_BG = RGBColor(0x00, 0x70, 0xA0)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF4, 0xF8)
SECTION_BG = RGBColor(0x00, 0x50, 0x80)

# Paths to outputs
BASE = "outputs/"


# ============================================================
# Helper functions
# ============================================================

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        if bold_first and ": " in item:
            bold_part, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = bold_part + ": "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return tf


def add_table(slide, left, top, width, height, data, col_widths=None,
              highlight_cols=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
                    if c >= 1:
                        paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
    return table


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
    else:
        print(f"  WARNING: Image not found: {path}")


def title_bar(slide, title, subtitle=None):
    add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), DARK_BG)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
                    subtitle, font_size=16, color=LIGHT_GRAY)


def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, SECTION_BG)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                    subtitle, font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4), Inches(3.55), Inches(5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()
    return slide


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
            "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e40\u0e27\u0e25\u0e32\u0e41\u0e1a\u0e1a\u0e41\u0e1b\u0e23\u0e1c\u0e31\u0e19",
            font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
            "\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e43\u0e19\u0e23\u0e30\u0e1a\u0e1a\u0e1c\u0e25\u0e34\u0e15\u0e44\u0e1f\u0e1f\u0e49\u0e32\u0e1e\u0e25\u0e31\u0e07\u0e07\u0e32\u0e19\u0e41\u0e2a\u0e07\u0e2d\u0e32\u0e17\u0e34\u0e15\u0e22\u0e4c",
            font_size=28, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(3.5), Inches(3.7), Inches(6), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.1), Inches(11), Inches(0.8),
            "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e02\u0e19\u0e32\u0e14 Input/Output: 4/1 vs 6/2 vs 8/3",
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
            "Site 1  |  \u0e21\u0e34\u0e16\u0e38\u0e19\u0e32\u0e22\u0e19 - \u0e18\u0e31\u0e19\u0e27\u0e32\u0e04\u0e21 2025  |  PyTorch  |  NVIDIA RTX 3060 Ti",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Experiment Setup
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e01\u0e32\u0e23\u0e15\u0e31\u0e49\u0e07\u0e04\u0e48\u0e32\u0e01\u0e32\u0e23\u0e17\u0e14\u0e25\u0e2d\u0e07", "3 \u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01 \u0e43\u0e0a\u0e49 6 \u0e42\u0e21\u0e40\u0e14\u0e25\u0e15\u0e48\u0e2d\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
            "\u0e04\u0e33\u0e16\u0e32\u0e21\u0e27\u0e34\u0e08\u0e31\u0e22: \u0e01\u0e32\u0e23\u0e40\u0e1b\u0e25\u0e35\u0e48\u0e22\u0e19\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07 Input/Output \u0e2a\u0e48\u0e07\u0e1c\u0e25\u0e15\u0e48\u0e2d\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e43\u0e19\u0e01\u0e32\u0e23\u0e17\u0e33\u0e19\u0e32\u0e22\u0e41\u0e25\u0e30\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e44\u0e23?",
            font_size=18, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(2.0), [
    ["\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01", "Input Window", "Output Window", "\u0e23\u0e30\u0e22\u0e30\u0e40\u0e27\u0e25\u0e32 Input", "\u0e23\u0e30\u0e22\u0e30\u0e40\u0e27\u0e25\u0e32\u0e17\u0e33\u0e19\u0e32\u0e22", "Trial Name"],
    ["\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (Baseline)", "4 timesteps", "1 timestep", "1 \u0e0a\u0e31\u0e48\u0e27\u0e42\u0e21\u0e07", "15 \u0e19\u0e32\u0e17\u0e35", "run_1"],
    ["\u0e02\u0e22\u0e32\u0e22 (Extended)", "6 timesteps", "2 timesteps", "1.5 \u0e0a\u0e31\u0e48\u0e27\u0e42\u0e21\u0e07", "30 \u0e19\u0e32\u0e17\u0e35", "exp_6step"],
    ["\u0e23\u0e30\u0e22\u0e30\u0e44\u0e01\u0e25 (Long-range)", "8 timesteps", "3 timesteps", "2 \u0e0a\u0e31\u0e48\u0e27\u0e42\u0e21\u0e07", "45 \u0e19\u0e32\u0e17\u0e35", "exp_8step"],
])

add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
            "\u0e41\u0e19\u0e27\u0e04\u0e34\u0e14 Sliding Window", font_size=20, color=ACCENT_BLUE, bold=True)

add_shape_bg(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(2.0), WHITE)
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.8),
            "\u0e15\u0e31\u0e27\u0e2d\u0e22\u0e48\u0e32\u0e07: --input-steps 6 --output-steps 2\n\n"
            "Time:  [t-5] [t-4] [t-3] [t-2] [t-1] [t]   |   [t+1] [t+2]\n"
            "        |__________ INPUT (6 x 15\u0e19\u0e32\u0e17\u0e35) ________|   |_ OUTPUT _|\n"
            "                    1.5 \u0e0a\u0e31\u0e48\u0e27\u0e42\u0e21\u0e07                   30 \u0e19\u0e32\u0e17\u0e35\n\n"
            "Input \u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19 = \u0e1a\u0e23\u0e34\u0e1a\u0e17\u0e21\u0e32\u0e01\u0e02\u0e36\u0e49\u0e19  |  Output \u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19 = \u0e17\u0e33\u0e19\u0e32\u0e22\u0e22\u0e32\u0e01\u0e02\u0e36\u0e49\u0e19",
            font_size=14, color=DARK_TEXT, font_name="Consolas")

# ============================================================
# SLIDE 3: Common Settings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e01\u0e32\u0e23\u0e15\u0e31\u0e49\u0e07\u0e04\u0e48\u0e32\u0e17\u0e35\u0e48\u0e43\u0e0a\u0e49\u0e23\u0e48\u0e27\u0e21\u0e01\u0e31\u0e19\u0e17\u0e38\u0e01\u0e01\u0e32\u0e23\u0e17\u0e14\u0e25\u0e2d\u0e07")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4),
            "\u0e04\u0e48\u0e32\u0e01\u0e32\u0e23\u0e1d\u0e36\u0e01\u0e2a\u0e2d\u0e19 (\u0e40\u0e2b\u0e21\u0e37\u0e2d\u0e19\u0e01\u0e31\u0e19\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01)", font_size=20, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(3.0), [
    ["\u0e1e\u0e32\u0e23\u0e32\u0e21\u0e34\u0e40\u0e15\u0e2d\u0e23\u0e4c", "\u0e04\u0e48\u0e32"],
    ["Epochs", "100"],
    ["Batch Size", "32"],
    ["Learning Rate", "0.001"],
    ["Early Stopping", "Patience = 15"],
    ["LR Scheduler", "ReduceLROnPlateau (0.5)"],
    ["Loss Function", "MSE"],
    ["\u0e40\u0e01\u0e13\u0e11\u0e4c\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34", "3-sigma"],
    ["Seed", "42"],
])

add_textbox(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(0.4),
            "6 \u0e42\u0e21\u0e40\u0e14\u0e25\u0e17\u0e35\u0e48\u0e1d\u0e36\u0e01\u0e2a\u0e2d\u0e19\u0e15\u0e48\u0e2d\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01", font_size=20, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.0), [
    ["#", "\u0e42\u0e21\u0e40\u0e14\u0e25", "\u0e1b\u0e23\u0e30\u0e40\u0e20\u0e17"],
    ["1", "Isolation Forest", "ML, Unsupervised"],
    ["2", "Random Forest", "ML, Regression"],
    ["3", "LSTM", "DL, \u0e25\u0e33\u0e14\u0e31\u0e1a\u0e40\u0e27\u0e25\u0e32"],
    ["4", "CNN-LSTM", "DL, \u0e44\u0e2e\u0e1a\u0e23\u0e34\u0e14"],
    ["5", "LSTM Autoencoder", "DL, \u0e2a\u0e23\u0e49\u0e32\u0e07\u0e43\u0e2b\u0e21\u0e48"],
    ["6", "Transformer", "DL, Self-Attention"],
])

add_textbox(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
            "\u0e02\u0e49\u0e2d\u0e21\u0e39\u0e25: 15,013 \u0e15\u0e31\u0e27\u0e2d\u0e22\u0e48\u0e32\u0e07 | \u0e1d\u0e36\u0e01: 10,509 (70%) | \u0e17\u0e14\u0e2a\u0e2d\u0e1a: 4,504 (30%) | \u0e1f\u0e35\u0e40\u0e08\u0e2d\u0e23\u0e4c: \u0e04\u0e48\u0e32\u0e41\u0e2a\u0e07 + \u0e2d\u0e38\u0e13\u0e2b\u0e20\u0e39\u0e21\u0e34 -> \u0e01\u0e33\u0e25\u0e31\u0e07\u0e1c\u0e25\u0e34\u0e15\u0e44\u0e1f\u0e1f\u0e49\u0e32",
            font_size=16, color=DARK_TEXT, bold=True)

add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
            "\u0e2b\u0e21\u0e32\u0e22\u0e40\u0e2b\u0e15\u0e38: \u0e42\u0e21\u0e40\u0e14\u0e25 ML (Isolation Forest, Random Forest) \u0e44\u0e21\u0e48\u0e43\u0e0a\u0e49 Sliding Window "
            "\u0e14\u0e31\u0e07\u0e19\u0e31\u0e49\u0e19\u0e1c\u0e25\u0e25\u0e31\u0e1e\u0e18\u0e4c\u0e08\u0e30\u0e40\u0e2b\u0e21\u0e37\u0e2d\u0e19\u0e01\u0e31\u0e19\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01 "
            "\u0e21\u0e35\u0e40\u0e1e\u0e35\u0e22\u0e07\u0e42\u0e21\u0e40\u0e14\u0e25 DL (LSTM, CNN-LSTM, LSTM Autoencoder, Transformer) \u0e17\u0e35\u0e48\u0e44\u0e14\u0e49\u0e23\u0e31\u0e1a\u0e1c\u0e25\u0e01\u0e23\u0e30\u0e17\u0e1a",
            font_size=15, color=ACCENT_RED)


# ============================================================
# SECTION: Results
# ============================================================
section_slide("\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e01\u0e32\u0e23\u0e17\u0e33\u0e19\u0e32\u0e22", "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e02\u0e2d\u0e07\u0e42\u0e21\u0e40\u0e14\u0e25 DL \u0e17\u0e31\u0e49\u0e07 3 \u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01")


# ============================================================
# SLIDE 4: DL Performance Comparison Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e42\u0e21\u0e40\u0e14\u0e25 DL",
          "RMSE / MAE / R\u00b2 \u0e1a\u0e19\u0e0a\u0e38\u0e14\u0e17\u0e14\u0e2a\u0e2d\u0e1a (RMSE & MAE \u0e22\u0e34\u0e48\u0e07\u0e15\u0e48\u0e33\u0e22\u0e34\u0e48\u0e07\u0e14\u0e35, R\u00b2 \u0e22\u0e34\u0e48\u0e07\u0e2a\u0e39\u0e07\u0e22\u0e34\u0e48\u0e07\u0e14\u0e35)")

# RMSE comparison
add_textbox(slide, Inches(0.6), Inches(1.4), Inches(3.8), Inches(0.4),
            "RMSE (\u0e04\u0e48\u0e32\u0e04\u0e25\u0e32\u0e14\u0e40\u0e04\u0e25\u0e37\u0e48\u0e2d\u0e19\u0e01\u0e33\u0e25\u0e31\u0e07\u0e2a\u0e2d\u0e07)", font_size=17, color=ACCENT_RED, bold=True)
add_table(slide, Inches(0.6), Inches(1.8), Inches(3.8), Inches(2.0), [
    ["\u0e42\u0e21\u0e40\u0e14\u0e25", "4/1", "6/2", "8/3"],
    ["LSTM", "31.81", "36.05", "39.55"],
    ["CNN-LSTM", "32.49", "38.21", "41.80"],
    ["Transformer", "32.22", "35.29", "39.48"],
])

# MAE comparison
add_textbox(slide, Inches(4.8), Inches(1.4), Inches(3.8), Inches(0.4),
            "MAE (\u0e04\u0e48\u0e32\u0e04\u0e25\u0e32\u0e14\u0e40\u0e04\u0e25\u0e37\u0e2d\u0e19\u0e40\u0e09\u0e25\u0e35\u0e48\u0e22)", font_size=17, color=ACCENT_ORANGE, bold=True)
add_table(slide, Inches(4.8), Inches(1.8), Inches(3.8), Inches(2.0), [
    ["\u0e42\u0e21\u0e40\u0e14\u0e25", "4/1", "6/2", "8/3"],
    ["LSTM", "12.80", "17.05", "20.12"],
    ["CNN-LSTM", "14.34", "17.79", "21.26"],
    ["Transformer", "13.72", "16.51", "19.19"],
])

# R² comparison
add_textbox(slide, Inches(9.0), Inches(1.4), Inches(3.8), Inches(0.4),
            "R\u00b2 (\u0e2a\u0e31\u0e21\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e4c\u0e01\u0e32\u0e23\u0e15\u0e31\u0e14\u0e2a\u0e34\u0e19\u0e43\u0e08)", font_size=17, color=ACCENT_GREEN, bold=True)
add_table(slide, Inches(9.0), Inches(1.8), Inches(3.8), Inches(2.0), [
    ["\u0e42\u0e21\u0e40\u0e14\u0e25", "4/1", "6/2", "8/3"],
    ["LSTM", "0.9467", "0.9304", "0.9150"],
    ["CNN-LSTM", "0.9444", "0.9219", "0.9051"],
    ["Transformer", "0.9453", "0.9333", "0.9153"],
])

# Key observations
add_shape_bg(slide, Inches(0.6), Inches(4.1), Inches(12.2), Inches(3.0), WHITE)
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4),
            "\u0e02\u0e49\u0e2d\u0e2a\u0e31\u0e07\u0e40\u0e01\u0e15\u0e2b\u0e25\u0e31\u0e01", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(4.7), Inches(11.5), Inches(2.5), [
    "\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (4/1): \u0e43\u0e2b\u0e49\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e14\u0e35\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e17\u0e38\u0e01\u0e42\u0e21\u0e40\u0e14\u0e25 - \u0e01\u0e32\u0e23\u0e17\u0e33\u0e19\u0e32\u0e22\u0e25\u0e48\u0e27\u0e07\u0e2b\u0e19\u0e49\u0e32 1 \u0e02\u0e31\u0e49\u0e19\u0e07\u0e48\u0e32\u0e22\u0e01\u0e27\u0e48\u0e32 2 \u0e2b\u0e23\u0e37\u0e2d 3 \u0e02\u0e31\u0e49\u0e19",
    "Transformer: \u0e40\u0e1b\u0e47\u0e19\u0e42\u0e21\u0e40\u0e14\u0e25\u0e17\u0e35\u0e48\u0e14\u0e35\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e43\u0e19\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01 - RMSE \u0e15\u0e48\u0e33\u0e2a\u0e38\u0e14\u0e41\u0e25\u0e30 R\u00b2 \u0e2a\u0e39\u0e07\u0e2a\u0e38\u0e14\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e2a\u0e21\u0e48\u0e33\u0e40\u0e2a\u0e21\u0e2d",
    "\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e25\u0e14\u0e25\u0e07\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e04\u0e48\u0e2d\u0e22\u0e40\u0e1b\u0e47\u0e19: R\u00b2 \u0e25\u0e14\u0e08\u0e32\u0e01 ~0.94-0.95 (4/1) \u0e40\u0e1b\u0e47\u0e19 ~0.93 (6/2) \u0e41\u0e25\u0e30 ~0.91 (8/3)",
    "CNN-LSTM: \u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e25\u0e14\u0e40\u0e23\u0e47\u0e27\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e40\u0e21\u0e37\u0e48\u0e2d\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19 - CNN Filter \u0e40\u0e2b\u0e21\u0e32\u0e30\u0e01\u0e31\u0e1a\u0e23\u0e39\u0e1b\u0e41\u0e1a\u0e1a\u0e2a\u0e31\u0e49\u0e19\u0e46",
    "\u0e42\u0e21\u0e40\u0e14\u0e25 DL \u0e17\u0e38\u0e01\u0e15\u0e31\u0e27\u0e22\u0e31\u0e07\u0e04\u0e07 R\u00b2 > 0.90 \u0e41\u0e21\u0e49\u0e41\u0e15\u0e48\u0e43\u0e19\u0e07\u0e32\u0e19\u0e17\u0e35\u0e48\u0e22\u0e32\u0e01\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14 (8-in/3-out) \u0e41\u0e2a\u0e14\u0e07\u0e27\u0e48\u0e32\u0e2a\u0e16\u0e32\u0e1b\u0e31\u0e15\u0e22\u0e01\u0e23\u0e23\u0e21\u0e41\u0e02\u0e47\u0e07\u0e41\u0e23\u0e07",
], font_size=15, bold_first=True)


# ============================================================
# SLIDE 5: Performance Degradation Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e27\u0e34\u0e40\u0e04\u0e23\u0e32\u0e30\u0e2b\u0e4c\u0e01\u0e32\u0e23\u0e25\u0e14\u0e25\u0e07\u0e02\u0e2d\u0e07\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e",
          "\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e25\u0e14\u0e25\u0e07\u0e40\u0e17\u0e48\u0e32\u0e44\u0e2b\u0e23\u0e48\u0e40\u0e21\u0e37\u0e48\u0e2d\u0e02\u0e22\u0e32\u0e22\u0e0a\u0e48\u0e27\u0e07\u0e01\u0e32\u0e23\u0e17\u0e33\u0e19\u0e32\u0e22?")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
            "RMSE \u0e17\u0e35\u0e48\u0e40\u0e1e\u0e34\u0e48\u0e21\u0e02\u0e36\u0e49\u0e19\u0e08\u0e32\u0e01\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (4/1)", font_size=20, color=ACCENT_RED, bold=True)

add_table(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(2.2), [
    ["\u0e42\u0e21\u0e40\u0e14\u0e25", "RMSE \u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19", "RMSE 6/2", "\u0394 \u0e08\u0e32\u0e01\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19", "RMSE 8/3", "\u0394 \u0e08\u0e32\u0e01\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19", "\u0394 \u0e23\u0e27\u0e21"],
    ["LSTM", "31.81", "36.05", "+4.24 (+13.3%)", "39.55", "+7.74 (+24.3%)", "+24.3%"],
    ["CNN-LSTM", "32.49", "38.21", "+5.72 (+17.6%)", "41.80", "+9.31 (+28.6%)", "+28.6%"],
    ["Transformer", "32.22", "35.29", "+3.07 (+9.5%)", "39.48", "+7.26 (+22.5%)", "+22.5%"],
])

add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
            "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a R\u00b2", font_size=20, color=ACCENT_GREEN, bold=True)

add_table(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8), [
    ["\u0e42\u0e21\u0e40\u0e14\u0e25", "R\u00b2 4/1", "R\u00b2 6/2", "\u0394 R\u00b2", "R\u00b2 8/3", "\u0394 R\u00b2", "\u0e2d\u0e31\u0e19\u0e14\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e41\u0e02\u0e47\u0e07\u0e41\u0e23\u0e07"],
    ["LSTM", "0.9467", "0.9304", "-0.0163", "0.9150", "-0.0317", "\u0e17\u0e35\u0e48 2"],
    ["CNN-LSTM", "0.9444", "0.9219", "-0.0225", "0.9051", "-0.0393", "\u0e17\u0e35\u0e48 3"],
    ["Transformer", "0.9453", "0.9333", "-0.0120", "0.9153", "-0.0300", "\u0e17\u0e35\u0e48 1 (\u0e41\u0e02\u0e47\u0e07\u0e41\u0e23\u0e07\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14)"],
])


# ============================================================
# SLIDE 6: Model Comparison Plots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e01\u0e23\u0e32\u0e1f\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e42\u0e21\u0e40\u0e14\u0e25")

add_textbox(slide, Inches(0.6), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (In=4 / Out=1)", font_size=15, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/Comparison/model_comparison.png",
               Inches(0.2), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(4.7), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e02\u0e22\u0e32\u0e22 (In=6 / Out=2)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/Comparison/model_comparison.png",
               Inches(4.5), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(8.9), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e23\u0e30\u0e22\u0e30\u0e44\u0e01\u0e25 (In=8 / Out=3)", font_size=15, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/Comparison/model_comparison.png",
               Inches(8.8), Inches(1.5), width=Inches(4.2))

add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(1.5), [
    "Random Forest RMSE/MAE \u0e04\u0e07\u0e17\u0e35\u0e48 (83.37 / 34.80) \u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01 - \u0e44\u0e21\u0e48\u0e43\u0e0a\u0e49 Sliding Window",
    "\u0e42\u0e21\u0e40\u0e14\u0e25 DL \u0e41\u0e2a\u0e14\u0e07 RMSE/MAE \u0e40\u0e1e\u0e34\u0e48\u0e21\u0e02\u0e36\u0e49\u0e19\u0e41\u0e25\u0e30 R\u00b2 \u0e25\u0e14\u0e25\u0e07\u0e40\u0e21\u0e37\u0e48\u0e2d\u0e0a\u0e48\u0e27\u0e07\u0e01\u0e32\u0e23\u0e17\u0e33\u0e19\u0e32\u0e22\u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19",
    "\u0e0a\u0e48\u0e2d\u0e07\u0e27\u0e48\u0e32\u0e07\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e23\u0e30\u0e2b\u0e27\u0e48\u0e32\u0e07 DL \u0e01\u0e31\u0e1a ML \u0e01\u0e27\u0e49\u0e32\u0e07\u0e02\u0e36\u0e49\u0e19\u0e40\u0e21\u0e37\u0e48\u0e2d\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e17\u0e33\u0e19\u0e32\u0e22\u0e2a\u0e31\u0e49\u0e19 (DL \u0e42\u0e14\u0e14\u0e40\u0e14\u0e48\u0e19\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e17\u0e35\u0e48 4/1)",
], font_size=15)


# ============================================================
# SECTION: Anomaly Detection
# ============================================================
section_slide("\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34",
              "\u0e01\u0e32\u0e23\u0e15\u0e31\u0e49\u0e07\u0e04\u0e48\u0e32 Timestep \u0e2a\u0e48\u0e07\u0e1c\u0e25\u0e15\u0e48\u0e2d\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e44\u0e23")


# ============================================================
# SLIDE 7: Anomaly Detection Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34: \u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e23\u0e32\u0e22\u0e42\u0e21\u0e40\u0e14\u0e25")

add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(3.2), [
    ["\u0e42\u0e21\u0e40\u0e14\u0e25", "\u0e27\u0e34\u0e18\u0e35", "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 4/1", "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 6/2", "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 8/3", "\u0e41\u0e19\u0e27\u0e42\u0e19\u0e49\u0e21"],
    ["Isolation Forest", "Isolation score", "244 (5.4%)", "244 (5.4%)", "244 (5.4%)", "\u0e04\u0e07\u0e17\u0e35\u0e48 (\u0e44\u0e21\u0e48\u0e43\u0e0a\u0e49 window)"],
    ["Random Forest", "3-sigma MAE", "159 (3.5%)", "159 (3.5%)", "159 (3.5%)", "\u0e04\u0e07\u0e17\u0e35\u0e48 (\u0e44\u0e21\u0e48\u0e43\u0e0a\u0e49 window)"],
    ["LSTM", "3-sigma MAE", "103 (2.3%)", "99 (2.2%)", "100 (2.2%)", "\u0e04\u0e07\u0e17\u0e35\u0e48 (~100)"],
    ["CNN-LSTM", "3-sigma MAE", "100 (2.2%)", "101 (2.2%)", "105 (2.3%)", "\u0e04\u0e07\u0e17\u0e35\u0e48 (~100)"],
    ["LSTM Autoencoder", "3-sigma recon.", "360 (8.0%)", "225 (5.0%)", "215 (4.8%)", "\u0e25\u0e14\u0e25\u0e07"],
    ["Transformer", "3-sigma MAE", "98 (2.2%)", "97 (2.2%)", "97 (2.2%)", "\u0e04\u0e07\u0e17\u0e35\u0e48\u0e21\u0e32\u0e01"],
])

add_shape_bg(slide, Inches(0.6), Inches(5.0), Inches(12.2), Inches(2.0), WHITE)
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.4),
            "\u0e02\u0e49\u0e2d\u0e04\u0e49\u0e19\u0e1e\u0e1a\u0e2b\u0e25\u0e31\u0e01", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(5.5), Inches(11.5), Inches(1.5), [
    "\u0e42\u0e21\u0e40\u0e14\u0e25 DL \u0e17\u0e33\u0e19\u0e32\u0e22 (LSTM, CNN-LSTM, Transformer): \u0e08\u0e33\u0e19\u0e27\u0e19\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e04\u0e07\u0e17\u0e35\u0e48\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e19\u0e48\u0e32\u0e17\u0e36\u0e48\u0e07\u0e43\u0e08\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01 (~97-105)",
    "LSTM Autoencoder: \u0e08\u0e33\u0e19\u0e27\u0e19\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e25\u0e14\u0e25\u0e07\u0e21\u0e32\u0e01 (360 -> 225 -> 215) - Input \u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19\u0e0a\u0e48\u0e27\u0e22\u0e40\u0e23\u0e35\u0e22\u0e19\u0e23\u0e39\u0e49\u0e23\u0e39\u0e1b\u0e41\u0e1a\u0e1a\u0e1b\u0e01\u0e15\u0e34\u0e44\u0e14\u0e49\u0e14\u0e35\u0e02\u0e36\u0e49\u0e19",
    "\u0e42\u0e21\u0e40\u0e14\u0e25 ML \u0e44\u0e21\u0e48\u0e40\u0e1b\u0e25\u0e35\u0e48\u0e22\u0e19\u0e41\u0e1b\u0e25\u0e07 (\u0e44\u0e21\u0e48\u0e02\u0e36\u0e49\u0e19\u0e01\u0e31\u0e1a\u0e25\u0e33\u0e14\u0e31\u0e1a): Isolation Forest=244, Random Forest=159 \u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01",
    "\u0e01\u0e0e 3-sigma \u0e1b\u0e23\u0e31\u0e1a\u0e40\u0e01\u0e13\u0e11\u0e4c\u0e42\u0e14\u0e22\u0e2d\u0e31\u0e15\u0e42\u0e19\u0e21\u0e31\u0e15\u0e34\u0e40\u0e21\u0e37\u0e48\u0e2d\u0e01\u0e32\u0e23\u0e01\u0e23\u0e30\u0e08\u0e32\u0e22\u0e15\u0e31\u0e27\u0e02\u0e2d\u0e07\u0e04\u0e48\u0e32\u0e04\u0e25\u0e32\u0e14\u0e40\u0e04\u0e25\u0e37\u0e2d\u0e19\u0e40\u0e1b\u0e25\u0e35\u0e48\u0e22\u0e19\u0e44\u0e1b\u0e15\u0e32\u0e21\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07",
], font_size=15, bold_first=True)


# ============================================================
# SLIDE 8: Anomaly Comparison Plots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e01\u0e23\u0e32\u0e1f\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e08\u0e33\u0e19\u0e27\u0e19\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34")

add_textbox(slide, Inches(0.6), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (In=4 / Out=1)", font_size=15, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/Comparison/anomaly_comparison.png",
               Inches(0.2), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(4.7), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e02\u0e22\u0e32\u0e22 (In=6 / Out=2)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/Comparison/anomaly_comparison.png",
               Inches(4.5), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(8.9), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e23\u0e30\u0e22\u0e30\u0e44\u0e01\u0e25 (In=8 / Out=3)", font_size=15, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/Comparison/anomaly_comparison.png",
               Inches(8.8), Inches(1.5), width=Inches(4.2))

add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(1.5), [
    "\u0e41\u0e17\u0e48\u0e07 LSTM Autoencoder \u0e2b\u0e14\u0e25\u0e07\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e40\u0e2b\u0e47\u0e19\u0e44\u0e14\u0e49\u0e0a\u0e31\u0e14\u0e08\u0e32\u0e01\u0e0b\u0e49\u0e32\u0e22\u0e44\u0e1b\u0e02\u0e27\u0e32 (360 -> 225 -> 215)",
    "\u0e42\u0e21\u0e40\u0e14\u0e25\u0e2d\u0e37\u0e48\u0e19\u0e17\u0e31\u0e49\u0e07\u0e2b\u0e21\u0e14\u0e21\u0e35\u0e08\u0e33\u0e19\u0e27\u0e19\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e04\u0e07\u0e17\u0e35\u0e48\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01",
    "\u0e04\u0e27\u0e32\u0e21\u0e04\u0e07\u0e17\u0e35\u0e48\u0e19\u0e35\u0e49\u0e40\u0e1b\u0e47\u0e19\u0e2a\u0e31\u0e0d\u0e0d\u0e32\u0e13\u0e17\u0e35\u0e48\u0e14\u0e35 - \u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e44\u0e21\u0e48\u0e44\u0e27\u0e15\u0e48\u0e2d\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e21\u0e32\u0e01\u0e40\u0e01\u0e34\u0e19\u0e44\u0e1b",
], font_size=15)


# ============================================================
# SLIDE 9: Ensemble Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e41\u0e1a\u0e1a Ensemble",
          "\u0e42\u0e2b\u0e27\u0e15\u0e40\u0e2a\u0e35\u0e22\u0e07\u0e02\u0e49\u0e32\u0e07\u0e21\u0e32\u0e01: \u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e40\u0e21\u0e37\u0e48\u0e2d >= 3 \u0e08\u0e32\u0e01 6 \u0e42\u0e21\u0e40\u0e14\u0e25\u0e40\u0e2b\u0e47\u0e19\u0e15\u0e23\u0e07\u0e01\u0e31\u0e19")

add_table(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.5), [
    ["\u0e23\u0e30\u0e14\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e40\u0e2b\u0e47\u0e19\u0e1e\u0e49\u0e2d\u0e07", "4/1", "6/2", "8/3"],
    ["0 \u0e42\u0e21\u0e40\u0e14\u0e25 (\u0e1b\u0e01\u0e15\u0e34)", "3,913 (86.9%)", "3,984 (88.5%)", "3,969 (89.3%)"],
    ["1 \u0e42\u0e21\u0e40\u0e14\u0e25", "289 (6.4%)", "259 (5.8%)", "277 (6.2%)"],
    ["2 \u0e42\u0e21\u0e40\u0e14\u0e25", "150 (3.3%)", "125 (2.8%)", "133 (3.0%)"],
    ["3 \u0e42\u0e21\u0e40\u0e14\u0e25", "129 (2.9%)", "109 (2.4%)", "94 (2.1%)"],
    ["4 \u0e42\u0e21\u0e40\u0e14\u0e25", "9 (0.2%)", "13 (0.3%)", "12 (0.3%)"],
    ["5 \u0e42\u0e21\u0e40\u0e14\u0e25", "8 (0.2%)", "5 (0.1%)", "7 (0.2%)"],
    ["6 \u0e42\u0e21\u0e40\u0e14\u0e25", "2 (0.04%)", "2 (0.04%)", "2 (0.04%)"],
])

add_shape_bg(slide, Inches(0.6), Inches(5.3), Inches(5.5), Inches(1.5), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(5), Inches(1.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 Ensemble \u0e23\u0e27\u0e21:\n"
            "4/1: 148 (3.3%)   |   6/2: 129 (2.9%)   |   8/3: 115 (2.6%)\n"
            "\u0e41\u0e19\u0e27\u0e42\u0e19\u0e49\u0e21\u0e25\u0e14\u0e25\u0e07\u0e40\u0e21\u0e37\u0e48\u0e2d\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19",
            font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Ensemble plots
add_textbox(slide, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.3),
            "\u0e01\u0e23\u0e32\u0e1f Ensemble", font_size=17, color=ACCENT_BLUE, bold=True)

add_textbox(slide, Inches(6.5), Inches(1.7), Inches(2.1), Inches(0.3),
            "4/1", font_size=13, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/Ensemble/plots/ensemble_anomalies.png",
               Inches(6.5), Inches(1.95), width=Inches(2.1))

add_textbox(slide, Inches(8.7), Inches(1.7), Inches(2.1), Inches(0.3),
            "6/2", font_size=13, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/Ensemble/plots/ensemble_anomalies.png",
               Inches(8.7), Inches(1.95), width=Inches(2.1))

add_textbox(slide, Inches(10.9), Inches(1.7), Inches(2.1), Inches(0.3),
            "8/3", font_size=13, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/Ensemble/plots/ensemble_anomalies.png",
               Inches(10.9), Inches(1.95), width=Inches(2.1))

add_bullet_list(slide, Inches(6.7), Inches(5.3), Inches(6), Inches(1.8), [
    "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 Ensemble \u0e25\u0e14\u0e25\u0e07: 148 -> 129 -> 115",
    "\u0e40\u0e01\u0e34\u0e14\u0e08\u0e32\u0e01 Autoencoder \u0e2a\u0e23\u0e49\u0e32\u0e07\u0e43\u0e2b\u0e21\u0e48\u0e44\u0e14\u0e49\u0e14\u0e35\u0e02\u0e36\u0e49\u0e19\u0e40\u0e21\u0e37\u0e48\u0e2d Input \u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19",
    "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e04\u0e27\u0e32\u0e21\u0e40\u0e0a\u0e37\u0e48\u0e2d\u0e21\u0e31\u0e48\u0e19\u0e2a\u0e39\u0e07 (6/6 \u0e42\u0e21\u0e40\u0e14\u0e25\u0e40\u0e2b\u0e47\u0e19\u0e15\u0e23\u0e07\u0e01\u0e31\u0e19) \u0e04\u0e07\u0e17\u0e35\u0e48 2 \u0e08\u0e38\u0e14",
], font_size=15)


# ============================================================
# SLIDE 10: Transformer Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e40\u0e08\u0e32\u0e30\u0e25\u0e36\u0e01\u0e42\u0e21\u0e40\u0e14\u0e25\u0e17\u0e35\u0e48\u0e14\u0e35\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14: Transformer \u0e17\u0e31\u0e49\u0e07 3 \u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01")

add_textbox(slide, Inches(0.6), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (In=4 / Out=1)", font_size=15, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/Transformer/plots/prediction_vs_actual.png",
               Inches(0.2), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(4.7), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e02\u0e22\u0e32\u0e22 (In=6 / Out=2)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/Transformer/plots/prediction_vs_actual.png",
               Inches(4.5), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(8.9), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e23\u0e30\u0e22\u0e30\u0e44\u0e01\u0e25 (In=8 / Out=3)", font_size=15, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/Transformer/plots/prediction_vs_actual.png",
               Inches(8.8), Inches(1.5), width=Inches(4.2))

# Transformer anomaly detection
add_textbox(slide, Inches(0.6), Inches(4.1), Inches(4), Inches(0.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 (4/1)", font_size=14, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/Transformer/plots/anomaly_detection.png",
               Inches(0.2), Inches(4.35), width=Inches(4.2))

add_textbox(slide, Inches(4.7), Inches(4.1), Inches(4), Inches(0.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 (6/2)", font_size=14, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/Transformer/plots/anomaly_detection.png",
               Inches(4.5), Inches(4.35), width=Inches(4.2))

add_textbox(slide, Inches(8.9), Inches(4.1), Inches(4), Inches(0.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 (8/3)", font_size=14, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/Transformer/plots/anomaly_detection.png",
               Inches(8.8), Inches(4.35), width=Inches(4.2))


# ============================================================
# SLIDE 11: LSTM Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a LSTM \u0e17\u0e31\u0e49\u0e07 3 \u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01")

add_textbox(slide, Inches(0.6), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (In=4 / Out=1)", font_size=15, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/LSTM/plots/prediction_vs_actual.png",
               Inches(0.2), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(4.7), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e02\u0e22\u0e32\u0e22 (In=6 / Out=2)", font_size=15, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/LSTM/plots/prediction_vs_actual.png",
               Inches(4.5), Inches(1.5), width=Inches(4.2))

add_textbox(slide, Inches(8.9), Inches(1.25), Inches(4), Inches(0.3),
            "\u0e23\u0e30\u0e22\u0e30\u0e44\u0e01\u0e25 (In=8 / Out=3)", font_size=15, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/LSTM/plots/prediction_vs_actual.png",
               Inches(8.8), Inches(1.5), width=Inches(4.2))

# LSTM anomaly detection
add_textbox(slide, Inches(0.6), Inches(4.1), Inches(4), Inches(0.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 (4/1)", font_size=14, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, BASE + "run_1/LSTM/plots/anomaly_detection.png",
               Inches(0.2), Inches(4.35), width=Inches(4.2))

add_textbox(slide, Inches(4.7), Inches(4.1), Inches(4), Inches(0.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 (6/2)", font_size=14, color=ACCENT_ORANGE, bold=True)
add_image_safe(slide, BASE + "exp_6step/LSTM/plots/anomaly_detection.png",
               Inches(4.5), Inches(4.35), width=Inches(4.2))

add_textbox(slide, Inches(8.9), Inches(4.1), Inches(4), Inches(0.3),
            "\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 (8/3)", font_size=14, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, BASE + "exp_8step/LSTM/plots/anomaly_detection.png",
               Inches(8.8), Inches(4.35), width=Inches(4.2))


# ============================================================
# SLIDE 12: Summary & Recommendation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e2a\u0e23\u0e38\u0e1b\u0e41\u0e25\u0e30\u0e02\u0e49\u0e2d\u0e41\u0e19\u0e30\u0e19\u0e33")

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(0.4),
            "\u0e2a\u0e23\u0e38\u0e1b\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e", font_size=22, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(0.6), Inches(2.0), Inches(6), Inches(2.0), [
    ["\u0e15\u0e31\u0e27\u0e27\u0e31\u0e14", "4/1 (\u0e14\u0e35\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14)", "6/2", "8/3"],
    ["RMSE \u0e15\u0e48\u0e33\u0e2a\u0e38\u0e14", "31.81 (LSTM)", "35.29 (Trans.)", "39.48 (Trans.)"],
    ["R\u00b2 \u0e2a\u0e39\u0e07\u0e2a\u0e38\u0e14", "0.9467 (LSTM)", "0.9333 (Trans.)", "0.9153 (Trans.)"],
    ["\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34 Ensemble", "148 (3.3%)", "129 (2.9%)", "115 (2.6%)"],
])

add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.8), Inches(0.4),
            "\u0e02\u0e49\u0e2d\u0e41\u0e19\u0e30\u0e19\u0e33\u0e15\u0e32\u0e21\u0e01\u0e23\u0e13\u0e35\u0e43\u0e0a\u0e49\u0e07\u0e32\u0e19", font_size=22, color=ACCENT_ORANGE, bold=True)

add_shape_bg(slide, Inches(7.2), Inches(2.0), Inches(5.6), Inches(1.2), WHITE)
add_textbox(slide, Inches(7.4), Inches(2.05), Inches(5.2), Inches(0.3),
            "\u0e15\u0e34\u0e14\u0e15\u0e32\u0e21\u0e41\u0e1a\u0e1a Real-time (\u0e41\u0e08\u0e49\u0e07\u0e40\u0e15\u0e37\u0e2d\u0e19\u0e17\u0e38\u0e01 15 \u0e19\u0e32\u0e17\u0e35)", font_size=16, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(7.4), Inches(2.35), Inches(5.2), Inches(0.7),
            "\u0e43\u0e0a\u0e49 In=4 / Out=1 - \u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e2a\u0e39\u0e07\u0e2a\u0e38\u0e14 \u0e17\u0e33\u0e19\u0e32\u0e22\u0e40\u0e23\u0e47\u0e27\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14",
            font_size=15, color=DARK_TEXT)

add_shape_bg(slide, Inches(7.2), Inches(3.3), Inches(5.6), Inches(1.2), WHITE)
add_textbox(slide, Inches(7.4), Inches(3.35), Inches(5.2), Inches(0.3),
            "\u0e1e\u0e22\u0e32\u0e01\u0e23\u0e13\u0e4c\u0e23\u0e30\u0e22\u0e30\u0e2a\u0e31\u0e49\u0e19 (\u0e25\u0e48\u0e27\u0e07\u0e2b\u0e19\u0e49\u0e32 30 \u0e19\u0e32\u0e17\u0e35)", font_size=16, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(7.4), Inches(3.65), Inches(5.2), Inches(0.7),
            "\u0e43\u0e0a\u0e49 In=6 / Out=2 - \u0e2a\u0e21\u0e14\u0e38\u0e25\u0e23\u0e30\u0e2b\u0e27\u0e48\u0e32\u0e07\u0e1a\u0e23\u0e34\u0e1a\u0e17\u0e41\u0e25\u0e30\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33 (R\u00b2=0.93)",
            font_size=15, color=DARK_TEXT)

add_shape_bg(slide, Inches(7.2), Inches(4.6), Inches(5.6), Inches(1.2), WHITE)
add_textbox(slide, Inches(7.4), Inches(4.65), Inches(5.2), Inches(0.3),
            "\u0e27\u0e32\u0e07\u0e41\u0e1c\u0e19\u0e01\u0e32\u0e23\u0e1c\u0e25\u0e34\u0e15 (\u0e25\u0e48\u0e27\u0e07\u0e2b\u0e19\u0e49\u0e32 45 \u0e19\u0e32\u0e17\u0e35)", font_size=16, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.4), Inches(4.95), Inches(5.2), Inches(0.7),
            "\u0e43\u0e0a\u0e49 In=8 / Out=3 - \u0e22\u0e31\u0e07\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e14\u0e35 (R\u00b2=0.91) \u0e2a\u0e33\u0e2b\u0e23\u0e31\u0e1a\u0e01\u0e32\u0e23\u0e27\u0e32\u0e07\u0e41\u0e1c\u0e19\u0e01\u0e32\u0e23\u0e40\u0e14\u0e34\u0e19\u0e40\u0e04\u0e23\u0e37\u0e48\u0e2d\u0e07",
            font_size=15, color=DARK_TEXT)

add_textbox(slide, Inches(0.6), Inches(4.3), Inches(6), Inches(0.4),
            "\u0e1b\u0e23\u0e30\u0e40\u0e14\u0e47\u0e19\u0e2a\u0e33\u0e04\u0e31\u0e0d", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(6), Inches(2.5), [
    "1. \u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e17\u0e33\u0e19\u0e32\u0e22\u0e2a\u0e31\u0e49\u0e19: \u0e43\u0e2b\u0e49\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e14\u0e35\u0e01\u0e27\u0e48\u0e32",
    "2. Transformer: \u0e41\u0e02\u0e47\u0e07\u0e41\u0e23\u0e07\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01",
    "3. \u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34: \u0e04\u0e07\u0e17\u0e35\u0e48\u0e44\u0e21\u0e48\u0e27\u0e48\u0e32\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07",
    "4. \u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01\u0e44\u0e14\u0e49 R\u00b2 > 0.90 - \u0e1e\u0e23\u0e49\u0e2d\u0e21\u0e43\u0e0a\u0e49\u0e07\u0e32\u0e19\u0e08\u0e23\u0e34\u0e07",
    "5. \u0e40\u0e25\u0e37\u0e2d\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01\u0e15\u0e32\u0e21\u0e04\u0e27\u0e32\u0e21\u0e15\u0e49\u0e2d\u0e07\u0e01\u0e32\u0e23\u0e43\u0e0a\u0e49\u0e07\u0e32\u0e19\u0e08\u0e23\u0e34\u0e07",
], font_size=16, bold_first=True)


# ============================================================
# SLIDE 13: Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SUBTLE_BG)
title_bar(slide, "\u0e2a\u0e23\u0e38\u0e1b\u0e1c\u0e25")

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(5.5), [
    "1. \u0e1e\u0e37\u0e49\u0e19\u0e10\u0e32\u0e19 (In=4/Out=1): \u0e43\u0e2b\u0e49\u0e04\u0e27\u0e32\u0e21\u0e41\u0e21\u0e48\u0e19\u0e22\u0e33\u0e14\u0e35\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14 Transformer RMSE=32.22, R\u00b2=0.9453 \u0e41\u0e25\u0e30 LSTM RMSE=31.81, R\u00b2=0.9467",
    "2. \u0e01\u0e32\u0e23\u0e02\u0e22\u0e32\u0e22\u0e0a\u0e48\u0e27\u0e07\u0e17\u0e33\u0e19\u0e32\u0e22: \u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e25\u0e14\u0e25\u0e07\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e04\u0e48\u0e2d\u0e22\u0e40\u0e1b\u0e47\u0e19 RMSE \u0e40\u0e1e\u0e34\u0e48\u0e21\u0e02\u0e36\u0e49\u0e19 ~22-29% \u0e08\u0e32\u0e01 4/1 \u0e44\u0e1b 8/3 \u0e41\u0e15\u0e48\u0e17\u0e38\u0e01\u0e42\u0e21\u0e40\u0e14\u0e25\u0e22\u0e31\u0e07\u0e04\u0e07 R\u00b2 > 0.90",
    "3. Transformer: \u0e41\u0e02\u0e47\u0e07\u0e41\u0e23\u0e07\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e17\u0e38\u0e01\u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01 RMSE \u0e40\u0e1e\u0e34\u0e48\u0e21\u0e19\u0e49\u0e2d\u0e22\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14 (+22.5%) \u0e41\u0e25\u0e30\u0e1b\u0e23\u0e30\u0e2a\u0e34\u0e17\u0e18\u0e34\u0e20\u0e32\u0e1e\u0e14\u0e35\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e2b\u0e23\u0e37\u0e2d\u0e40\u0e01\u0e37\u0e2d\u0e1a\u0e17\u0e35\u0e48\u0e2a\u0e38\u0e14\u0e2d\u0e22\u0e48\u0e32\u0e07\u0e2a\u0e21\u0e48\u0e33\u0e40\u0e2a\u0e21\u0e2d",
    "4. \u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e04\u0e07\u0e17\u0e35\u0e48: \u0e42\u0e21\u0e40\u0e14\u0e25 DL \u0e17\u0e33\u0e19\u0e32\u0e22\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e44\u0e14\u0e49 ~97-105 \u0e08\u0e38\u0e14\u0e44\u0e21\u0e48\u0e27\u0e48\u0e32\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e08\u0e30\u0e40\u0e1b\u0e47\u0e19\u0e40\u0e17\u0e48\u0e32\u0e44\u0e2b\u0e23\u0e48",
    "5. LSTM Autoencoder: \u0e44\u0e14\u0e49\u0e1b\u0e23\u0e30\u0e42\u0e22\u0e0a\u0e19\u0e4c\u0e08\u0e32\u0e01 Input \u0e22\u0e32\u0e27\u0e02\u0e36\u0e49\u0e19 \u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e25\u0e14\u0e08\u0e32\u0e01 360 \u0e40\u0e1b\u0e47\u0e19 215 \u0e41\u0e2a\u0e14\u0e07\u0e27\u0e48\u0e32\u0e40\u0e23\u0e35\u0e22\u0e19\u0e23\u0e39\u0e49\u0e23\u0e39\u0e1b\u0e41\u0e1a\u0e1a\u0e1b\u0e01\u0e15\u0e34\u0e44\u0e14\u0e49\u0e14\u0e35\u0e02\u0e36\u0e49\u0e19",
    "6. Ensemble Voting: \u0e1b\u0e23\u0e31\u0e1a\u0e15\u0e31\u0e27\u0e42\u0e14\u0e22\u0e2d\u0e31\u0e15\u0e42\u0e19\u0e21\u0e31\u0e15\u0e34 \u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e25\u0e14 (148 -> 129 -> 115) \u0e41\u0e15\u0e48\u0e04\u0e27\u0e32\u0e21\u0e40\u0e0a\u0e37\u0e48\u0e2d\u0e21\u0e31\u0e48\u0e19\u0e2a\u0e39\u0e07 (6/6 \u0e40\u0e2b\u0e47\u0e19\u0e15\u0e23\u0e07\u0e01\u0e31\u0e19) \u0e04\u0e07\u0e17\u0e35\u0e48 2 \u0e08\u0e38\u0e14",
    "7. \u0e2a\u0e16\u0e32\u0e1b\u0e31\u0e15\u0e22\u0e01\u0e23\u0e23\u0e21 Variable Timestep: \u0e1c\u0e48\u0e32\u0e19\u0e01\u0e32\u0e23\u0e17\u0e14\u0e2a\u0e2d\u0e1a\u0e41\u0e25\u0e49\u0e27 \u0e42\u0e04\u0e49\u0e14\u0e40\u0e14\u0e35\u0e22\u0e27\u0e23\u0e2d\u0e07\u0e23\u0e31\u0e1a\u0e2b\u0e25\u0e32\u0e22\u0e42\u0e2b\u0e21\u0e14\u0e01\u0e32\u0e23\u0e17\u0e33\u0e07\u0e32\u0e19\u0e42\u0e14\u0e22\u0e44\u0e21\u0e48\u0e15\u0e49\u0e2d\u0e07\u0e40\u0e1b\u0e25\u0e35\u0e48\u0e22\u0e19\u0e2a\u0e16\u0e32\u0e1b\u0e31\u0e15\u0e22\u0e01\u0e23\u0e23\u0e21",
], font_size=17, bold_first=True)


# ============================================================
# SLIDE 14: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
            "\u0e02\u0e2d\u0e1a\u0e04\u0e38\u0e13\u0e04\u0e23\u0e31\u0e1a",
            font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4), Inches(3.3), Inches(5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_ORANGE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1.2),
            "\u0e40\u0e1b\u0e23\u0e35\u0e22\u0e1a\u0e40\u0e17\u0e35\u0e22\u0e1a\u0e02\u0e19\u0e32\u0e14\u0e2b\u0e19\u0e49\u0e32\u0e15\u0e48\u0e32\u0e07\u0e40\u0e27\u0e25\u0e32\u0e41\u0e1a\u0e1a\u0e41\u0e1b\u0e23\u0e1c\u0e31\u0e19\n\u0e01\u0e32\u0e23\u0e15\u0e23\u0e27\u0e08\u0e08\u0e31\u0e1a\u0e04\u0e27\u0e32\u0e21\u0e1c\u0e34\u0e14\u0e1b\u0e01\u0e15\u0e34\u0e43\u0e19\u0e23\u0e30\u0e1a\u0e1a\u0e1c\u0e25\u0e34\u0e15\u0e44\u0e1f\u0e1f\u0e49\u0e32\u0e1e\u0e25\u0e31\u0e07\u0e07\u0e32\u0e19\u0e41\u0e2a\u0e07\u0e2d\u0e32\u0e17\u0e34\u0e15\u0e22\u0e4c\n\u0e14\u0e49\u0e27\u0e22 Machine Learning \u0e41\u0e25\u0e30 Deep Learning",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "3 \u0e04\u0e2d\u0e19\u0e1f\u0e34\u0e01  |  6 \u0e42\u0e21\u0e40\u0e14\u0e25  |  18 \u0e23\u0e2d\u0e1a\u0e01\u0e32\u0e23\u0e1d\u0e36\u0e01  |  Site 1  |  PyTorch",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = "documents/Solar_Anomaly_Detection_compare_TH.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
