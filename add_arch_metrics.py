from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ppt_path = "documents/run_training_summary.pptx"
prs = Presentation(ppt_path)

blank = prs.slide_layouts[5]
slide = prs.slides.add_slide(blank)
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
title_tf = title.text_frame
title_tf.text = "Model Architecture Overview"
title_tf.paragraphs[0].font.size = Pt(32)
title_tf.paragraphs[0].font.bold = True

def add_block(text, x, y):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(1.0))
    shape.text_frame.text = text
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    return shape

start = add_block("Raw data\n(datasets/site_1)", 0.5, 1.2)
proc = add_block("Preprocessing + scaling\n(MinMax, sequences)", 3.0, 1.2)
pred = add_block("Deep predictors\nLSTM / CNN-LSTM / Transformer\nout_steps x n_features", 5.5, 1.2)
ae = add_block("LSTM Autoencoder\n(reconstruction)", 5.5, 3.0)
out = add_block("Outputs + Anomaly flags\n3 sigma + ensemble", 3.0, 3.0)

def add_arrow(x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.width = Pt(2)

add_arrow(2.3, 1.7, 3.0, 1.7)
add_arrow(4.8, 1.7, 5.5, 1.7)
add_arrow(6.4, 1.7, 6.4, 3.0)
add_arrow(4.4, 2.7, 5.5, 2.7)
add_arrow(2.3, 2.7, 3.0, 2.7)

layout = prs.slide_layouts[1]
metrics_slide = prs.slides.add_slide(layout)
metrics_slide.shapes.title.text = "Metric Summary Template"
rows, cols = 5, 5
left, top, width, height = Inches(0.4), Inches(1.4), Inches(9.2), Inches(3)
table = metrics_slide.shapes.add_table(rows, cols, left, top, width, height).table
headers = ["Model", "MAE", "RMSE", "Anomaly Count", "Notes"]
for idx, title in enumerate(headers):
    cell = table.cell(0, idx)
    cell.text = title
    cell.text_frame.paragraphs[0].font.size = Pt(12)

rows_data = [
    ["LSTM", "TBD", "TBD", "TBD", "Default predictor"],
    ["CNN-LSTM", "TBD", "TBD", "TBD", "CNN preprocessing"],
    ["Transformer", "TBD", "TBD", "TBD", "Attention pooled"],
]
for r, data in enumerate(rows_data, start=1):
    for c, value in enumerate(data):
        cell = table.cell(r, c)
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(12)

prs.save(ppt_path)
print(ppt_path)
