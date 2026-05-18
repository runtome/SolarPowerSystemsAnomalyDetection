#!/usr/bin/env python3
"""
make_report.py  —  Generate a PPTX summary report from a training run output directory.

Usage:
    python make_report.py                                   # auto-find latest run_N
    python make_report.py --output-dir outputs/run_4
    python make_report.py --output-dir outputs/run_4 --output-pptx my_report.pptx
"""

import argparse
import json
import re
from datetime import datetime
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

TITLE_TOP    = Inches(0.20)
TITLE_HEIGHT = Inches(0.65)
CONTENT_TOP  = Inches(0.92)
CONTENT_H    = Inches(6.40)
MARGIN       = Inches(0.35)
CONTENT_W    = SLIDE_W - 2 * MARGIN

COLOR_DARK_BLUE  = RGBColor(0x1F, 0x49, 0x7D)
COLOR_ACCENT     = RGBColor(0x2E, 0x75, 0xB6)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK      = RGBColor(0x00, 0x00, 0x00)
COLOR_HEADER_BG  = RGBColor(0x2E, 0x75, 0xB6)
COLOR_ALT_ROW    = RGBColor(0xDD, 0xEB, 0xF7)
COLOR_SUBTITLE   = RGBColor(0xBD, 0xD7, 0xEE)
COLOR_GRAY_TEXT  = RGBColor(0x77, 0x77, 0x77)

DL_MODELS  = ["LSTM", "CNN_LSTM", "LSTM_Autoencoder", "Transformer"]
ALL_MODELS = ["Isolation_Forest", "Random_Forest"] + DL_MODELS


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def safe_img(path) -> Path | None:
    p = Path(path)
    return p if p.exists() else None


def img_size_px(path: Path):
    with Image.open(path) as im:
        return im.size  # (width_px, height_px)


def fit_image(slide, img_path: Path, left: Emu, top: Emu, max_w: Emu, max_h: Emu):
    """Add image centered inside (max_w × max_h), preserving aspect ratio."""
    w_px, h_px = img_size_px(img_path)
    ratio = w_px / h_px
    if max_w / max_h >= ratio:
        height = int(max_h)
        width  = int(max_h * ratio)
    else:
        width  = int(max_w)
        height = int(max_w / ratio)
    h_offset = (max_h - height) // 2
    w_offset = (max_w - width)  // 2
    slide.shapes.add_picture(str(img_path),
                              left + w_offset, top + h_offset,
                              width, height)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text: str, left, top, width, height,
                 font_size=14, bold=False, color=COLOR_BLACK,
                 align=PP_ALIGN.LEFT, bg_color=None, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def add_slide_title(slide, title: str, subtitle: str = ""):
    """Blue title bar spanning the full width."""
    add_rect(slide, 0, TITLE_TOP, SLIDE_W, TITLE_HEIGHT, COLOR_ACCENT)
    label = f"  {title}" + (f"   —   {subtitle}" if subtitle else "")
    add_text_box(slide, label,
                 Inches(0.1), TITLE_TOP, SLIDE_W - Inches(0.2), TITLE_HEIGHT,
                 font_size=20, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_cover_slide(prs, trial_name: str, sites: list):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_DARK_BLUE)
    add_rect(slide, 0, SLIDE_H - Inches(0.65), SLIDE_W, Inches(0.65), COLOR_ACCENT)

    add_text_box(slide, "Solar Power Systems Anomaly Detection",
                 MARGIN, Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(1.2),
                 font_size=38, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    today = datetime.now().strftime("%Y-%m-%d")
    sub = f"Trial: {trial_name}   |   Sites: {', '.join(sites)}   |   {today}"
    add_text_box(slide, sub,
                 MARGIN, Inches(3.1), SLIDE_W - 2 * MARGIN, Inches(0.7),
                 font_size=16, color=COLOR_SUBTITLE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Machine Learning & Deep Learning Anomaly Detection Pipeline",
                 MARGIN, SLIDE_H - Inches(0.60), SLIDE_W - 2 * MARGIN, Inches(0.55),
                 font_size=12, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


def add_section_slide(prs, title: str, subtitle: str = ""):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_DARK_BLUE)
    add_text_box(slide, title,
                 MARGIN, Inches(2.7), SLIDE_W - 2 * MARGIN, Inches(1.1),
                 font_size=42, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, subtitle,
                     MARGIN, Inches(3.95), SLIDE_W - 2 * MARGIN, Inches(0.65),
                     font_size=18, color=COLOR_SUBTITLE, align=PP_ALIGN.CENTER)


def add_kv_table_slide(prs, title: str, rows: list):
    """Two-column key / value table."""
    slide = blank_slide(prs)
    add_slide_title(slide, title)

    n = len(rows)
    tbl_w = Inches(9)
    tbl_h = min(Inches(5.8), Emu(int(Inches(0.45)) * (n + 1)))
    tbl_left = (SLIDE_W - tbl_w) // 2
    tbl_top  = CONTENT_TOP + Inches(0.1)

    table = slide.shapes.add_table(n + 1, 2, tbl_left, tbl_top, tbl_w, tbl_h).table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(5)

    for ci, hdr in enumerate(["Parameter", "Value"]):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.bold  = True
        run.font.size  = Pt(13)
        run.font.color.rgb = COLOR_WHITE

    for ri, (k, v) in enumerate(rows):
        for ci, val in enumerate([k, str(v)]):
            cell = table.cell(ri + 1, ci)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ALT_ROW
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(12)
            if ci == 0:
                run.font.bold = True


def add_metrics_table_slide(prs, title: str, headers: list, rows: list,
                              bold_row_idx: int = None, subtitle: str = ""):
    slide = blank_slide(prs)
    add_slide_title(slide, title, subtitle)

    n_cols = len(headers)
    tbl_w  = CONTENT_W
    tbl_h  = min(Inches(5.5), Emu(int(Inches(0.44)) * (len(rows) + 1)))

    table = slide.shapes.add_table(
        len(rows) + 1, n_cols,
        MARGIN, CONTENT_TOP + Inches(0.1),
        tbl_w, tbl_h
    ).table

    col_w = tbl_w // n_cols
    for ci in range(n_cols):
        table.columns[ci].width = col_w

    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold  = True
        run.font.size  = Pt(12)
        run.font.color.rgb = COLOR_WHITE

    for ri, row_data in enumerate(rows):
        is_best = (ri == bold_row_idx)
        for ci, val in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ALT_ROW
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            if is_best:
                run.font.bold = True
                run.font.color.rgb = COLOR_ACCENT


def _place_img_or_placeholder(slide, img_path, left, top, w, h, label=""):
    if img_path and safe_img(img_path):
        fit_image(slide, Path(img_path), left, top, w, h)
    else:
        name = Path(img_path).name if img_path else "N/A"
        add_text_box(slide, f"[not found]\n{name}",
                     left, top, w, h,
                     font_size=9, color=COLOR_GRAY_TEXT,
                     align=PP_ALIGN.CENTER)


def add_two_image_slide(prs, title: str, left_img, right_img,
                         left_label: str = "", right_label: str = ""):
    slide = blank_slide(prs)
    add_slide_title(slide, title)

    label_h = Inches(0.28) if (left_label or right_label) else Inches(0)
    img_h   = CONTENT_H - label_h
    half_w  = (CONTENT_W - Inches(0.15)) // 2
    gap     = Inches(0.15)

    for img_path, lft, label in [
        (left_img,  MARGIN,               left_label),
        (right_img, MARGIN + half_w + gap, right_label),
    ]:
        _place_img_or_placeholder(slide, img_path, lft, CONTENT_TOP, half_w, img_h, label)
        if label:
            add_text_box(slide, label,
                         lft, CONTENT_TOP + img_h, half_w, label_h,
                         font_size=10, align=PP_ALIGN.CENTER, color=COLOR_GRAY_TEXT)


def add_single_image_slide(prs, title: str, img_path, subtitle: str = ""):
    slide = blank_slide(prs)
    add_slide_title(slide, title, subtitle)
    _place_img_or_placeholder(slide, img_path, MARGIN, CONTENT_TOP, CONTENT_W, CONTENT_H)


def add_four_image_slide(prs, title: str, img_paths: list, labels: list):
    """2 × 2 image grid."""
    slide = blank_slide(prs)
    add_slide_title(slide, title)

    gap    = Inches(0.15)
    half_w = (CONTENT_W - gap) // 2
    label_h = Inches(0.26)
    half_h  = (CONTENT_H - gap) // 2
    img_h   = half_h - label_h

    positions = [
        (MARGIN,           CONTENT_TOP),
        (MARGIN + half_w + gap, CONTENT_TOP),
        (MARGIN,           CONTENT_TOP + half_h + gap),
        (MARGIN + half_w + gap, CONTENT_TOP + half_h + gap),
    ]

    for i, (img_path, label) in enumerate(zip(img_paths, labels)):
        lft, tp = positions[i]
        _place_img_or_placeholder(slide, img_path, lft, tp, half_w, img_h)
        add_text_box(slide, label,
                     lft, tp + img_h, half_w, label_h,
                     font_size=10, align=PP_ALIGN.CENTER, color=COLOR_GRAY_TEXT)


# ---------------------------------------------------------------------------
# Data loaders
# ---------------------------------------------------------------------------

def load_run_config(output_dir: Path, sites: list) -> dict:
    for site in sites:
        p = output_dir / "Comparison" / site / "run_config.json"
        if p.exists():
            return json.loads(p.read_text())
    return {}


def load_model_metrics(output_dir: Path, site: str) -> pd.DataFrame | None:
    p = output_dir / "Comparison" / site / "model_comparison.csv"
    return pd.read_csv(p) if p.exists() else None


def load_anomaly_counts(output_dir: Path, site: str) -> dict:
    p = output_dir / "Comparison" / site / "anomaly_counts.json"
    return json.loads(p.read_text()) if p.exists() else {}


def load_ensemble_results(output_dir: Path, site: str) -> dict:
    p = output_dir / "Ensemble" / site / "plots" / "ensemble_results.json"
    return json.loads(p.read_text()) if p.exists() else {}


# ---------------------------------------------------------------------------
# Report builder
# ---------------------------------------------------------------------------

def build_report(output_dir: Path, output_pptx: Path):
    sites = sorted([d.name for d in (output_dir / "Comparison").iterdir() if d.is_dir()])
    if not sites:
        raise FileNotFoundError(f"No sites found under {output_dir}/Comparison/")

    trial_name = output_dir.name
    run_cfg    = load_run_config(output_dir, sites)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 01: Cover ─────────────────────────────────────────────────────
    add_cover_slide(prs, trial_name, sites)

    # ── Slide 02: Run Configuration ─────────────────────────────────────────
    cfg_rows = [
        ("Trial Name",    run_cfg.get("trial_name", trial_name)),
        ("Input Steps",   f"{run_cfg.get('input_steps', '—')} steps = "
                          f"{run_cfg.get('input_steps', 0) * 15} min"),
        ("Output Steps",  f"{run_cfg.get('output_steps', '—')} step = "
                          f"{run_cfg.get('output_steps', 0) * 15} min"),
        ("Epochs",        run_cfg.get("epochs", "—")),
        ("Batch Size",    run_cfg.get("batch_size", "—")),
        ("Learning Rate", run_cfg.get("lr", "—")),
        ("Anomaly Sigma", run_cfg.get("sigma", "—")),
        ("Train Ratio",   f"{run_cfg.get('train_ratio', '—')} (70/30 time-based split)"),
        ("Features",      ", ".join(run_cfg.get("features", []))),
        ("Target",        run_cfg.get("target", "generation_kw")),
        ("Sites",         ", ".join(sites)),
    ]
    add_kv_table_slide(prs, "Run Configuration", cfg_rows)

    # ── Slide 03: Cross-site summary ─────────────────────────────────────────
    summary_rows = []
    for site in sites:
        df  = load_model_metrics(output_dir, site)
        ens = load_ensemble_results(output_dir, site)
        if df is not None and not df.empty:
            best = df.loc[df["R2"].idxmax()]
            summary_rows.append([
                site,
                best["model"],
                f"{best['RMSE']:.2f}",
                f"{best['MAE']:.2f}",
                f"{best['R2']:.4f}",
                str(ens.get("total_anomalies", "—")),
            ])
        else:
            summary_rows.append([site, "—", "—", "—", "—",
                                  str(ens.get("total_anomalies", "—"))])

    add_metrics_table_slide(prs,
        "Cross-Site Performance Summary",
        ["Site", "Best Model", "RMSE", "MAE", "R²", "Ensemble Anomalies"],
        summary_rows,
        subtitle="Best model = highest R² on test set")

    # ── Per-site slides ──────────────────────────────────────────────────────
    for site in sites:
        metrics_df     = load_model_metrics(output_dir, site)
        anomaly_counts = load_anomaly_counts(output_dir, site)
        ens_results    = load_ensemble_results(output_dir, site)

        eda_dir = output_dir / "data_EDA" / site
        cmp_dir = output_dir / "Comparison"  / site
        ens_dir = output_dir / "Ensemble"    / site / "plots"

        n_anomalies = ens_results.get("total_anomalies", "—")

        # Section header
        add_section_slide(prs, f"Site: {site}",
                          f"Ensemble anomalies: {n_anomalies}")

        # EDA 1: time series + anomaly indicators
        add_two_image_slide(prs,
            "EDA — Time Series & Anomaly Indicators",
            eda_dir / "03_time_series.png",
            eda_dir / "12_anomaly_indicators.png",
            "Time Series Overview",
            "Anomaly Indicators Over Time")

        # EDA 2: distributions + correlation
        add_two_image_slide(prs,
            "EDA — Feature Distributions & Correlation",
            eda_dir / "05_distributions.png",
            eda_dir / "07_correlation_matrix.png",
            "Feature Distributions",
            "Correlation Matrix")

        # Model metrics table
        if metrics_df is not None and not metrics_df.empty:
            best_idx = int(metrics_df["R2"].idxmax())
            tbl_rows = []
            for _, row in metrics_df.iterrows():
                model = row["model"]
                ac = anomaly_counts.get(model, "—")
                tbl_rows.append([
                    model,
                    f"{row['RMSE']:.2f}",
                    f"{row['MAE']:.2f}",
                    f"{row['R2']:.4f}",
                    str(ac),
                ])
            # Append unsupervised models missing from model_comparison.csv
            for m in ["Isolation_Forest", "LSTM_Autoencoder"]:
                if m not in metrics_df["model"].values and m in anomaly_counts:
                    tbl_rows.append([m, "—", "—", "—", str(anomaly_counts[m])])
            add_metrics_table_slide(prs,
                "Model Performance",
                ["Model", "RMSE", "MAE", "R²", "Anomalies"],
                tbl_rows,
                bold_row_idx=best_idx,
                subtitle=f"Site: {site}  |  Bold row = best R²")

        # Model comparison charts
        add_two_image_slide(prs,
            "Model Comparison Charts",
            cmp_dir / "model_comparison.png",
            cmp_dir / "train_test_split.png",
            "RMSE / MAE / R² per Model",
            "Train / Test Split Timeline")

        # Anomaly detection results
        add_two_image_slide(prs,
            "Anomaly Detection Results",
            cmp_dir / "anomaly_comparison.png",
            ens_dir / "ensemble_anomalies.png",
            "Anomaly Counts per Model",
            "Ensemble Anomalies on Timeline")

        # Ensemble categories
        add_single_image_slide(prs,
            "Ensemble — Anomaly Categories",
            ens_dir / "anomaly_categories.png",
            subtitle=f"Site: {site}")

        # DL model loss curves (2×2)
        loss_imgs   = [output_dir / m / site / "loss" / "training_loss.png" for m in DL_MODELS]
        loss_labels = [m.replace("_", " ") for m in DL_MODELS]
        add_four_image_slide(prs,
            f"DL Model Training Loss — {site}",
            loss_imgs, loss_labels)

    # ── Last slide: report info ─────────────────────────────────────────────
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_DARK_BLUE)
    add_text_box(slide, "Report generated by make_report.py",
                 MARGIN, Inches(2.5), SLIDE_W - 2 * MARGIN, Inches(0.65),
                 font_size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    gen_time = datetime.now().strftime("%Y-%m-%d %H:%M")
    info = (f"Trial: {trial_name}   |   Sites: {', '.join(sites)}\n"
            f"Models: {', '.join(ALL_MODELS)}\n"
            f"Generated: {gen_time}")
    add_text_box(slide, info,
                 MARGIN, Inches(3.3), SLIDE_W - 2 * MARGIN, Inches(1.5),
                 font_size=14, color=COLOR_SUBTITLE, align=PP_ALIGN.CENTER)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    print(f"\n  Saved  : {output_pptx}")
    print(f"  Slides : {len(prs.slides)}")
    print(f"  Sites  : {', '.join(sites)}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def auto_detect_output_dir(base: str = "outputs") -> Path:
    base_p = Path(base)
    candidates = sorted(
        [d for d in base_p.iterdir()
         if d.is_dir() and re.match(r"run_\d+$", d.name)],
        key=lambda d: int(re.search(r"\d+", d.name).group())
    )
    if not candidates:
        raise FileNotFoundError(f"No run_N directories found in '{base}/'")
    return candidates[-1]


def main():
    parser = argparse.ArgumentParser(
        description="Generate PPTX summary report from a training run output directory.")
    parser.add_argument("--output-dir", default=None,
                        help="Trial output directory (default: auto-detect latest run_N)")
    parser.add_argument("--output-pptx", default=None,
                        help="Output file path (default: {output_dir}/{trial}_summary_report.pptx)")
    args = parser.parse_args()

    if args.output_dir is None:
        output_dir = auto_detect_output_dir()
        print(f"Auto-detected output dir: {output_dir}")
    else:
        output_dir = Path(args.output_dir)

    if not output_dir.exists():
        raise FileNotFoundError(f"Output directory not found: {output_dir}")

    trial_name = output_dir.name
    output_pptx = Path(args.output_pptx) if args.output_pptx \
        else output_dir / f"{trial_name}_summary_report.pptx"

    print(f"Building report for: {output_dir}")
    build_report(output_dir, output_pptx)


if __name__ == "__main__":
    main()
