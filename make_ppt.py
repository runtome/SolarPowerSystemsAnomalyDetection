from pptx import Presentation
from pathlib import Path

prs = Presentation()
layout = prs.slide_layouts

slide = prs.slides.add_slide(layout[0])
slide.shapes.title.text = "Run Training Overview"
slide.placeholders[1].text = "Solar Power Anomaly Detection pipeline summary"

slide = prs.slides.add_slide(layout[1])
slide.shapes.title.text = "Data & Model scope"
body = slide.placeholders[1].text_frame
body.text = "Inputs & Target"
p = body.add_paragraph()
p.text = "- features: irradiance_wm2, temperature_c (resampled 15 min)"
p.level = 1
p = body.add_paragraph()
p.text = "- target: generation_kw with MinMax scaling"
p.level = 1
p = body.add_paragraph()
p.text = "- Sequence config defaults: input=4, output=1 (1h in -> 15 min out)"
p.level = 1
p = body.add_paragraph()
p.text = "DL models: LSTM, CNN-LSTM, Transformer (predict output_steps x n_features)"
p.level = 0
p = body.add_paragraph()
p.text = "Autoencoder reconstructs input window for anomaly detection"
p.level = 1
p = body.add_paragraph()
p.text = "ML models: Isolation Forest (contamination 2%) + Random Forest regressor"
p.level = 0

slide = prs.slides.add_slide(layout[1])
slide.shapes.title.text = "Training flow & outputs"
body = slide.placeholders[1].text_frame
body.text = "Training config"
p = body.add_paragraph()
p.text = "- epochs=100, batch=32, Adam lr=0.001, patience=15, ReduceLROnPlateau factor=0.5"
p.level = 1
p = body.add_paragraph()
p.text = "- Loss: MSE on flattened output_steps x n_features vector"
p.level = 1
p = body.add_paragraph()
p.text = "- Logs: train/val loss history JSON/CSV, saved checkpoints, plots per model"
p.level = 1
p = body.add_paragraph()
p.text = "Anomaly detection: reconstruction MAE or 3 sigma thresholding and ensemble votes"
p.level = 0
p = body.add_paragraph()
p.text = "Comparison outputs: run_config.json, model_comparison.csv, anomaly counts"
p.level = 1

slide = prs.slides.add_slide(layout[1])
slide.shapes.title.text = "Sequence variations"
body = slide.placeholders[1].text_frame
body.text = "Variants to run"
p = body.add_paragraph()
p.text = "1) default: --input-steps 4 --output-steps 1 (1h in -> 15 min out)"
p.level = 1
p = body.add_paragraph()
p.text = "2) --input-steps 6 --output-steps 2 (1.5h in -> 30 min out)"
p.level = 1
p = body.add_paragraph()
p.text = "3) --input-steps 8 --output-steps 3 (2h in -> 45 min out)"
p.level = 1
p = body.add_paragraph()
p.text = "DL outputs grow with output_steps x n_features; data module adjusts splits"
p.level = 1

slide = prs.slides.add_slide(layout[1])
slide.shapes.title.text = "Summary & next steps"
body = slide.placeholders[1].text_frame
body.text = "Recommendations"
p = body.add_paragraph()
p.text = "- Review loss/anomaly plots for each sequence variant"
p.level = 1
p = body.add_paragraph()
p.text = "- Compare ML vs DL anomaly counts; use ensemble votes if >=2 models"
p.level = 1
p = body.add_paragraph()
p.text = "- Tune lr/batch-size or expand features via config if needed"
p.level = 1

out_path = Path("documents/run_training_summary.pptx")
out_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(out_path)
print(out_path)
